"""Quick verification that the purple/gold alternating pattern was applied."""
from pptx import Presentation

prs = Presentation(r"C:\Users\mohan\Desktop\3.pptx")
slide = prs.slides[0]

box_owner = {
    "Rectangle 26":  "Box 1 accent",
    "Rounded Rectangle 27":  "Box 1 badge ",
    "Rectangle 43":  "Box 2 accent",
    "Rounded Rectangle 44":  "Box 2 badge ",
    "Rectangle 55":  "Box 3 accent",
    "Rounded Rectangle 56":  "Box 3 badge ",
    "Rectangle 78":  "Box 4 accent",
    "Rounded Rectangle 79":  "Box 4 badge ",
    "Rectangle 86":  "Box 5 accent",
    "Rounded Rectangle 87":  "Box 5 badge ",
    "Rectangle 95":  "Box 6 accent",
    "Rounded Rectangle 96":  "Box 6 badge ",
    "Rectangle 165": "Box 7 accent",  # Eval
    "Rectangle 203": "Box 8 accent",  # Conclusion
    "Rounded Rectangle 204": "Box 8 badge ",
    "Rectangle 115": "Box 9 accent",  # Future Work
}


def color_name(rgb_str):
    if rgb_str == "8B5BC4":
        return "purple"
    if rgb_str == "D49A1F":
        return "gold"
    return f"OTHER({rgb_str})"


print("--- Verification: alternating Purple(8B5BC4) / Gold(D49A1F) pattern ---")
for shape in slide.shapes:
    if shape.name in box_owner:
        try:
            rgb = str(shape.fill.fore_color.rgb)
            print(f"  {box_owner[shape.name]}: #{rgb}  ->  {color_name(rgb)}")
        except Exception as e:
            print(f"  {box_owner[shape.name]}: ERR {e}")

# Rounded Rectangle 166 appears TWICE (one Eval, one FW).
rr166 = sorted([s for s in slide.shapes if s.name == "Rounded Rectangle 166"],
               key=lambda s: s.left)
if len(rr166) >= 1:
    rgb = str(rr166[0].fill.fore_color.rgb)
    print(f"  Box 7 badge : #{rgb}  ->  {color_name(rgb)}")
if len(rr166) >= 2:
    rgb = str(rr166[1].fill.fore_color.rgb)
    print(f"  Box 9 badge : #{rgb}  ->  {color_name(rgb)}")
