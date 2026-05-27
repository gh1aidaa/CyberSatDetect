# -*- coding: utf-8 -*-
"""
Build the professional English CyberSatDetect poster (A0 portrait, 84.1 x 118.9 cm).

Design language (mirrors the reference infographic the user provided):
    - Dark navy/teal gradient background
    - Teal + gold accent palette
    - Cards with header bands and circular icons
    - Real project visuals (bar chart + threshold scatter from
      `thesis_official_evaluation_figures/`).
    - English copy sourced entirely from this repository's docs/eval files.

All static images come from scripts/poster_assets/ which is produced by
scripts/poster_assets/generate_assets.py.
"""
from __future__ import annotations

import io
import os
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Emu, Pt

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "scripts" / "poster_assets"

# ---------------------------------------------------------------------------
# Palette (must match generate_assets.py)
# ---------------------------------------------------------------------------
C_BG_TOP      = "08131F"
C_BG_MID      = "0F2638"
C_BG_BOTTOM   = "0A1A28"
C_PANEL       = "0F2A3A"
C_PANEL_EDGE  = "1F4A60"
C_PANEL_HDR   = "163B52"
C_ACCENT      = "3FB7C2"
C_ACCENT2     = "5DE0E6"
C_ACCENT3     = "86F2F0"
C_GOLD        = "E8B842"
C_WHITE       = "FFFFFF"
C_TEXT        = "D7E6EC"
C_MUTED       = "8FA6B0"
C_GREEN       = "3DD68C"
C_RED         = "E5675A"
C_DIVIDER     = "26506A"

POSTER_W_CM = 84.1
POSTER_H_CM = 118.9

# ---------------------------------------------------------------------------
# Low-level shape helpers
# ---------------------------------------------------------------------------

def add_picture(slide, path: Path, left_cm: float, top_cm: float,
                width_cm: float | None = None, height_cm: float | None = None):
    args = [str(path), Cm(left_cm), Cm(top_cm)]
    kwargs = {}
    if width_cm is not None:
        kwargs["width"] = Cm(width_cm)
    if height_cm is not None:
        kwargs["height"] = Cm(height_cm)
    return slide.shapes.add_picture(*args, **kwargs)


def add_rect(slide, left_cm: float, top_cm: float, width_cm: float,
             height_cm: float, *, fill_hex: str | None = None,
             outline_hex: str | None = None, outline_pt: float = 0,
             rounded: bool = False, transparency: int = 0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Cm(left_cm), Cm(top_cm),
                                  Cm(width_cm), Cm(height_cm))
    if rounded:
        shp.adjustments[0] = 0.05
    if fill_hex is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
        if transparency:
            _set_fill_transparency(shp, transparency)
    if outline_hex is None or outline_pt <= 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = RGBColor.from_string(outline_hex)
        shp.line.width = Pt(outline_pt)
    shp.shadow.inherit = False
    if shp.has_text_frame:
        shp.text_frame.margin_left = Cm(0)
        shp.text_frame.margin_right = Cm(0)
        shp.text_frame.margin_top = Cm(0)
        shp.text_frame.margin_bottom = Cm(0)
    return shp


def _set_fill_transparency(shape, pct: int) -> None:
    """Add a:alpha to the solid fill."""
    sp_pr = shape.fill._xPr  # noqa: SLF001
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    solid = sp_pr.find(f"{{{ns}}}solidFill")
    if solid is None:
        return
    srgb = solid.find(f"{{{ns}}}srgbClr")
    if srgb is None:
        return
    # Remove existing alpha then add
    for child in list(srgb):
        if etree.QName(child).localname == "alpha":
            srgb.remove(child)
    alpha = etree.SubElement(srgb, f"{{{ns}}}alpha")
    alpha.set("val", str(int((100 - pct) * 1000)))


def add_text(slide, left_cm: float, top_cm: float, width_cm: float,
             height_cm: float, *, lines: list[tuple[str, dict]] | list[str],
             color_hex: str = C_TEXT, size_pt: float = 14, bold: bool = False,
             font_name: str = "Calibri", align: str = "l",
             v_align: str = "t", line_spacing: float | None = None) -> object:
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm),
                                  Cm(height_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    tf.auto_size = None
    if v_align == "c":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_align == "b":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    align_map = {"l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT, "c": PP_ALIGN.CENTER,
                 "j": PP_ALIGN.JUSTIFY}

    def _apply_line(p, items):
        if isinstance(items, str):
            items = [(items, {})]
        first = True
        for chunk, style in items:
            if first:
                run = p.add_run()
                first = False
            else:
                run = p.add_run()
            run.text = chunk
            font = run.font
            font.name = style.get("font_name", font_name)
            font.size = Pt(style.get("size_pt", size_pt))
            font.bold = style.get("bold", bold)
            font.color.rgb = RGBColor.from_string(style.get(
                "color_hex", color_hex))

    for i, item in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if isinstance(item, str):
            _apply_line(p, [(item, {})])
        elif isinstance(item, tuple) and len(item) == 2 and isinstance(
                item[1], dict):
            _apply_line(p, [item])
        else:
            _apply_line(p, item)
    return tb


# ---------------------------------------------------------------------------
# Card builder (used for every section panel)
# ---------------------------------------------------------------------------

def add_card(slide, left, top, width, height, *,
             title: str, icon_path: Path | None = None,
             header_height: float = 3.2, body_padding: float = 1.0,
             header_color: str = C_PANEL_HDR) -> tuple[float, float]:
    # Panel background
    add_rect(slide, left, top, width, height,
             fill_hex=C_PANEL, outline_hex=C_PANEL_EDGE, outline_pt=1.5,
             rounded=True, transparency=10)
    # Header band - solid, slightly taller than text height
    add_rect(slide, left, top, width, header_height,
             fill_hex=header_color, outline_hex=C_ACCENT, outline_pt=1.5,
             rounded=True)
    # Gold accent stripe along left edge of header
    add_rect(slide, left + 0.6, top + 0.6, 0.45, header_height - 1.2,
             fill_hex=C_GOLD)
    # Icon (left side of header)
    icon_size = header_height - 1.2
    if icon_path is not None and icon_path.exists():
        add_picture(slide, icon_path, left + 1.4, top + 0.6,
                    width_cm=icon_size, height_cm=icon_size)
        title_x = left + 1.4 + icon_size + 0.7
    else:
        title_x = left + 1.4
    # Title
    add_text(slide, title_x, top + 0.4,
             width - (title_x - left) - 0.4, header_height - 0.8,
             lines=[title],
             color_hex=C_WHITE, size_pt=24, bold=True, align="l", v_align="c",
             font_name="Calibri")
    # Return body region top-left for further content
    return (left + body_padding, top + header_height + 0.6)


# ---------------------------------------------------------------------------
# Tagline builder (4 across the top below the title)
# ---------------------------------------------------------------------------

def add_tagline(slide, left, top, width, height, *,
                icon_path: Path, headline: str, sub: str) -> None:
    icon_size = height - 0.8
    add_picture(slide, icon_path, left, top + 0.4,
                width_cm=icon_size, height_cm=icon_size)
    text_x = left + icon_size + 0.8
    add_text(slide, text_x, top + 0.3, width - icon_size - 0.8, height - 0.6,
             lines=[
                 (headline, {"size_pt": 19, "bold": True,
                             "color_hex": C_WHITE}),
                 (sub, {"size_pt": 14, "color_hex": C_TEXT}),
             ], align="l", v_align="c", line_spacing=1.2)


# ---------------------------------------------------------------------------
# Content (sourced from this repo)
# ---------------------------------------------------------------------------

PROJECT_TITLE = "CyberSatDetect"
PROJECT_SUB   = "An AI-Driven Anomaly Detection System for Satellite Security"
PROJECT_TAG   = "Continual learning that defends what powers the world above."

TAGLINES = [
    ("icon_shield.png", "Strengthens space asset security",
     "Operational continuity for critical satellite systems."),
    ("icon_dish.png", "Analyzes telemetry remotely",
     "Detects deviations from normal in-orbit behavior."),
    ("icon_pulse.png", "Continually learns new threats",
     "Adapts to novel and unknown attack patterns."),
    ("icon_brain.png", "AI-powered intelligent defense",
     "Hybrid Autoencoder + Predictor detects unseen attacks."),
]

INTRO = [
    "Satellites are critical global infrastructure for communication, navigation, and Earth observation. Any disruption to their telemetry can compromise vital missions and put space assets at risk.",
    "Cyberattacks on telemetry use patterns (drift, freeze, noise, spike) that fixed-threshold rules often miss, requiring a system that learns normal behavior and flags deviations.",
    "CyberSatDetect is an unsupervised hybrid Autoencoder + Predictor trained on normal telemetry only, paired with a safe continual-learning module that adapts to new operating regimes without catastrophic forgetting.",
]

OBJECTIVES = [
    ("Build an unsupervised Hybrid Autoencoder + Predictor that learns from normal telemetry only.", "icon_brain.png"),
    ("Achieve high separation between normal behavior and four attack types (Drift, Freeze, Noise, Spike).", "icon_target.png"),
    ("Reduce false alarms via statistical thresholds derived from normal data only (p95, p97, p99, p99.5, p99.7, 3σ).", "icon_shield.png"),
    ("Enable safe continual learning that ingests new normal windows under strict admin approval and replay.", "icon_loop.png"),
    ("Provide operational governance with a model registry (PENDING → APPROVED) and one-click rollback.", "icon_gear.png"),
]

METHODOLOGY = [
    ("1", "Telemetry ingestion", "Accepts CSV and NPY with strict numeric cleaning (float32, NaN repair, linear interpolation)."),
    ("2", "Sliding windows", "W = 100 timesteps with stride S = 50 (50% overlap), formatted as X ∈ ℝ^(B×100×1)."),
    ("3", "Hybrid Autoencoder + Predictor", "Two heads: reconstruction + next-step prediction with composite loss W_recon·L_recon + W_pred·L_pred + W_grad·L_grad + W_sep·L_sep."),
    ("4", "Anomaly score", "score = e_recon + e_pred + e_grad (no separation term at inference)."),
    ("5", "Statistical thresholds", "p99 / p99.5 / p99.7 / 3σ computed from normal data only."),
    ("6", "Continual learning", "buffer_manager → admin approval → build_dataset → fine_tune → model_registry."),
]

RESULTS_BULLETS_LEFT = [
    ("F1-Score", "0.949"),
    ("Accuracy", "98.57%"),
    ("Balanced Acc.", "99.14%"),
    ("Recall", "0.9991"),
    ("Precision", "0.9043"),
    ("FAR", "1.63%"),
    ("ROC-AUC", "0.996"),
    ("PR-AUC", "0.956"),
]

RECOMMENDATIONS = [
    "Deploy CyberSatDetect at ground stations as a real-time monitoring layer integrated with SIEM / SOC pipelines.",
    "Expand attack coverage to compound patterns (Drift+Spike, Pattern Shift, Scale, Drop) for stronger generalization.",
    "Tune the operating threshold to the FAR budget: p99 to suppress alarms, best-F1 when missing detections is costly.",
    "Activate the continual-learning loop with strict human approval: build_dataset → fine_tune → registry → APPROVED.",
    "Store explainability records per anomalous window (e_recon, e_pred, e_grad components) for post-incident review.",
    "Run periodic re-evaluation on production data to detect distribution drift and recalibrate thresholds.",
]

CONCLUSION = [
    "CyberSatDetect proves that a hybrid Autoencoder + Predictor trained on normal telemetry only can detect four major satellite attack types with high fidelity (F1 = 0.949, ROC-AUC = 0.996, FAR ≈ 1.63%).",
    "Future Work: multi-channel telemetry • live operational evaluation • adaptive threshold policy that follows distribution drift • compound attack scenarios • multi-task continual learning.",
]

TEAM = [
    "Ghaidaa A. Algarni",
    "Alaa F. Alhazmi",
    "Bayan A. Alzahrani",
    "Ghaidaa A. Almasudi",
]
SUPERVISOR = "Dr. Ehad A. Aljarf"
COLLEGE = "College of Computing — Department of Cybersecurity — Umm Al-Qura University"

# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

def build() -> Path:
    prs = Presentation()
    prs.slide_width  = Cm(POSTER_W_CM)
    prs.slide_height = Cm(POSTER_H_CM)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # ---------------------------------------------------------------- 1) BG
    bg_path = ASSETS / "bg.png"
    if bg_path.exists():
        add_picture(slide, bg_path, 0, 0, width_cm=POSTER_W_CM,
                    height_cm=POSTER_H_CM)

    # ---------------------------------------------------------------- 2) Header
    header_top = 4.0
    header_h   = 18.0

    # Satellite illustration on the right
    sat = ASSETS / "satellite.png"
    if sat.exists():
        add_picture(slide, sat, POSTER_W_CM - 36.0, header_top - 1.0,
                    width_cm=36.0, height_cm=26.0)

    # Title block left — "CyberSatDetect:" massive headline
    add_text(slide, 4.5, header_top, 48, 9.2,
             lines=[(PROJECT_TITLE + ":", {"size_pt": 92, "bold": True,
                                            "color_hex": C_WHITE})],
             align="l", v_align="t", line_spacing=1.0)
    # Subtitle in two lines for impact
    add_text(slide, 4.5, header_top + 7.4, 48, 7.5,
             lines=[
                 ("An AI-Driven Anomaly Detection", {"size_pt": 38,
                                                       "bold": True,
                                                       "color_hex": C_ACCENT3}),
                 ("System for Satellite Security", {"size_pt": 38,
                                                       "bold": True,
                                                       "color_hex": C_ACCENT3}),
             ], align="l", line_spacing=1.05)
    # Gold accent underline
    add_rect(slide, 4.5, header_top + 15.0, 22, 0.30, fill_hex=C_GOLD)
    add_text(slide, 4.5, header_top + 15.6, 48, 2.0,
             lines=[PROJECT_TAG], color_hex=C_GOLD, size_pt=20, bold=True,
             align="l")

    # ---------------------------------------------------------------- 3) Taglines
    tag_top = header_top + header_h + 1.5
    tag_h   = 6.8
    tag_left = 4.5
    tag_width_total = POSTER_W_CM - 2 * 4.5
    tag_count = len(TAGLINES)
    gap = 1.2
    tag_w = (tag_width_total - gap * (tag_count - 1)) / tag_count

    # Backing divider line above taglines
    add_rect(slide, 4.5, tag_top - 0.8, tag_width_total, 0.08,
             fill_hex=C_DIVIDER)

    for i, (icon, head, sub) in enumerate(TAGLINES):
        x = tag_left + i * (tag_w + gap)
        add_tagline(slide, x, tag_top, tag_w, tag_h,
                    icon_path=ASSETS / icon, headline=head, sub=sub)

    # ---------------------------------------------------------------- 4) Cards
    content_top = tag_top + tag_h + 2.2
    content_bottom_limit = POSTER_H_CM - 12.0
    content_h = content_bottom_limit - content_top

    col_gap = 1.8
    col_w   = (POSTER_W_CM - 2 * 4.5 - col_gap) / 2
    left_x  = 4.5
    right_x = 4.5 + col_w + col_gap

    row_gap = 1.5
    row_h   = (content_h - 2 * row_gap) / 3

    # Row 1 left — Introduction
    bx, by = add_card(slide, left_x, content_top, col_w, row_h,
                      title="INTRODUCTION",
                      icon_path=ASSETS / "icon_alert.png")
    intro_text_h = row_h - 12.0
    add_text(slide, bx, by, col_w - 2 * 1.0, intro_text_h,
             lines=[(p, {"size_pt": 17}) for p in INTRO],
             color_hex=C_TEXT, line_spacing=1.3, align="j")

    # Key facts strip at the bottom of the Introduction card
    facts = [
        ("100%", "Unsupervised\n(normal-only training)"),
        ("4", "Attack types\n(Drift • Freeze • Noise • Spike)"),
        ("W=100", "Window size with\n50% overlap stride"),
        ("Hybrid", "Autoencoder + Predictor\ncomposite loss"),
    ]
    facts_top = by + intro_text_h + 0.6
    fact_gap = 0.3
    fact_w = (col_w - 2.0 - fact_gap * (len(facts) - 1)) / len(facts)
    fact_h = 5.5
    for i, (val, lbl) in enumerate(facts):
        fx = bx + i * (fact_w + fact_gap)
        add_rect(slide, fx, facts_top, fact_w, fact_h,
                 fill_hex=C_PANEL_HDR, outline_hex=C_ACCENT, outline_pt=1.0,
                 rounded=True)
        add_text(slide, fx, facts_top + 0.4, fact_w, 2.0,
                 lines=[(val, {"size_pt": 26, "bold": True,
                                "color_hex": C_GOLD})],
                 align="c", v_align="c")
        add_text(slide, fx + 0.2, facts_top + 2.4, fact_w - 0.4, 2.7,
                 lines=[(lbl, {"size_pt": 11, "color_hex": C_TEXT})],
                 align="c", v_align="t", line_spacing=1.2)

    # Row 1 right — Objectives
    bx2, by2 = add_card(slide, right_x, content_top, col_w, row_h,
                        title="OBJECTIVES",
                        icon_path=ASSETS / "icon_target.png")
    obj_y = by2
    bullet_pad = 0.4
    for text, icon_name in OBJECTIVES:
        icon_path = ASSETS / icon_name
        if icon_path.exists():
            add_picture(slide, icon_path, bx2, obj_y, width_cm=1.6,
                        height_cm=1.6)
        add_text(slide, bx2 + 2.1, obj_y + 0.1, col_w - 3.2, 2.4,
                 lines=[(text, {"size_pt": 16})], color_hex=C_TEXT,
                 line_spacing=1.2)
        obj_y += 2.4 + bullet_pad

    # Row 2 left — Methodology
    row2_top = content_top + row_h + row_gap
    bx, by = add_card(slide, left_x, row2_top, col_w, row_h,
                      title="METHODOLOGY  &  TOOLS",
                      icon_path=ASSETS / "icon_gear.png")
    method_y = by
    for num, head, body in METHODOLOGY:
        # numbered circle badge
        badge = add_rect(slide, bx, method_y + 0.1, 1.4, 1.4,
                          fill_hex=C_GOLD, rounded=True)
        badge.adjustments[0] = 0.5
        add_text(slide, bx, method_y + 0.1, 1.4, 1.4,
                 lines=[(num, {"size_pt": 18, "bold": True,
                                "color_hex": "0F2638"})],
                 align="c", v_align="c")
        add_text(slide, bx + 2.0, method_y, col_w - 3.2, 1.5,
                 lines=[(head, {"size_pt": 16, "bold": True,
                                 "color_hex": C_ACCENT3})], align="l")
        add_text(slide, bx + 2.0, method_y + 1.5, col_w - 3.2, 1.5,
                 lines=[(body, {"size_pt": 14, "color_hex": C_TEXT})],
                 line_spacing=1.2)
        method_y += 2.65
    # Tool chips
    chips_y = method_y + 0.1
    chips = ["Python", "TensorFlow / Keras", "NumPy", "FastAPI", "React",
             "Docker"]
    chip_x = bx
    for c in chips:
        approx_w = 0.32 * len(c) + 1.4
        add_rect(slide, chip_x, chips_y, approx_w, 1.3,
                 fill_hex=C_PANEL_HDR, outline_hex=C_ACCENT, outline_pt=1.0,
                 rounded=True)
        add_text(slide, chip_x, chips_y, approx_w, 1.3,
                 lines=[(c, {"size_pt": 13, "bold": True,
                              "color_hex": C_ACCENT3})],
                 align="c", v_align="c")
        chip_x += approx_w + 0.4

    # Row 2 right — Results & Discussion
    bx2, by2 = add_card(slide, right_x, row2_top, col_w, row_h,
                        title="RESULTS  &  DISCUSSION",
                        icon_path=ASSETS / "icon_chart.png")

    # Left half: metric cards (2x4 grid)
    metrics_w = (col_w - 2.0) * 0.42
    metrics_x = bx2
    metric_card_w = metrics_w / 2 - 0.2
    metric_card_h = 2.1
    for i, (label, value) in enumerate(RESULTS_BULLETS_LEFT):
        r = i // 2
        c = i % 2
        mx = metrics_x + c * (metric_card_w + 0.4)
        my = by2 + r * (metric_card_h + 0.3)
        add_rect(slide, mx, my, metric_card_w, metric_card_h,
                 fill_hex=C_PANEL_HDR, outline_hex=C_ACCENT, outline_pt=1.0,
                 rounded=True)
        add_text(slide, mx, my + 0.15, metric_card_w, 1.0,
                 lines=[(value, {"size_pt": 22, "bold": True,
                                  "color_hex": C_GOLD})],
                 align="c", v_align="c")
        add_text(slide, mx, my + 1.15, metric_card_w, 0.9,
                 lines=[(label, {"size_pt": 12, "color_hex": C_TEXT})],
                 align="c", v_align="c")

    # Right half: bar chart
    chart_x = metrics_x + metrics_w + 0.6
    chart_w = col_w - 2.0 - metrics_w - 0.6
    chart_h = row_h - 5.0
    bar_path = ASSETS / "bar_chart.png"
    if bar_path.exists():
        add_picture(slide, bar_path, chart_x, by2,
                    width_cm=chart_w, height_cm=chart_h)

    # Test-set caption under results
    add_text(slide, bx2, by2 + (metric_card_h + 0.3) * 4 + 0.3,
             col_w - 2.0, 2.0,
             lines=[
                 ("Test set: ", {"size_pt": 13, "bold": True,
                                  "color_hex": C_ACCENT3}),
                 ("72 files • 226,623 normal windows • 182,323 attacked windows • 54,712 anomalous (drift, freeze, noise, spike).",
                  {"size_pt": 13, "color_hex": C_TEXT}),
             ], align="l", line_spacing=1.2)

    # Row 3 left — Recommendations
    row3_top = row2_top + row_h + row_gap
    bx, by = add_card(slide, left_x, row3_top, col_w, row_h,
                      title="RECOMMENDATIONS",
                      icon_path=ASSETS / "icon_dish.png")
    ry = by
    for line in RECOMMENDATIONS:
        # gold dot bullet
        dot = add_rect(slide, bx, ry + 0.4, 0.5, 0.5, fill_hex=C_GOLD,
                       rounded=True)
        dot.adjustments[0] = 0.5
        add_text(slide, bx + 1.0, ry, col_w - 2.2, 2.2,
                 lines=[(line, {"size_pt": 15, "color_hex": C_TEXT})],
                 line_spacing=1.2)
        ry += 2.45

    # Row 3 right — Conclusion & Future Work
    bx2, by2 = add_card(slide, right_x, row3_top, col_w, row_h,
                        title="CONCLUSION  &  FUTURE WORK",
                        icon_path=ASSETS / "icon_shield.png")
    # Big highlight box for the headline finding
    highlight_h = 7.0
    add_rect(slide, bx2, by2, col_w - 2.0, highlight_h,
             fill_hex=C_PANEL_HDR, outline_hex=C_GOLD, outline_pt=1.8,
             rounded=True)
    # Gold-bordered "key result" badge
    add_rect(slide, bx2 + 0.4, by2 + 0.4, 4.5, 1.4,
             fill_hex=C_GOLD, rounded=True)
    add_text(slide, bx2 + 0.4, by2 + 0.4, 4.5, 1.4,
             lines=[("KEY RESULT", {"size_pt": 13, "bold": True,
                                      "color_hex": "0F2638"})],
             align="c", v_align="c")
    add_text(slide, bx2 + 0.6, by2 + 2.0, col_w - 3.2, highlight_h - 2.3,
             lines=[(CONCLUSION[0], {"size_pt": 17, "color_hex": C_WHITE})],
             align="j", line_spacing=1.3)

    # Future work block below highlight
    fw_top = by2 + highlight_h + 1.2
    add_text(slide, bx2, fw_top, col_w - 2.0, 1.8,
             lines=[("FUTURE WORK", {"size_pt": 18, "bold": True,
                                       "color_hex": C_ACCENT3})], align="l")
    # Bullet items for future work
    future_bullets = [
        "Multi-channel telemetry support beyond the current single-channel scope.",
        "Live operational evaluation on real ground-station traffic.",
        "Adaptive threshold policy that tracks distribution drift automatically.",
        "Compound attack scenarios (Drift+Spike, Pattern Shift, Scale, Drop).",
        "Multi-task continual learning across diverse satellite operating regimes.",
    ]
    fy = fw_top + 2.0
    for line in future_bullets:
        dot = add_rect(slide, bx2 + 0.1, fy + 0.4, 0.45, 0.45,
                       fill_hex=C_GOLD, rounded=True)
        dot.adjustments[0] = 0.5
        add_text(slide, bx2 + 0.9, fy, col_w - 2.5, 1.5,
                 lines=[(line, {"size_pt": 14, "color_hex": C_TEXT})],
                 line_spacing=1.2)
        fy += 1.55

    # ---------------------------------------------------------------- 5) Footer
    footer_top = POSTER_H_CM - 10.5
    # Top divider line
    add_rect(slide, 4.5, footer_top - 0.5, POSTER_W_CM - 9.0, 0.08,
             fill_hex=C_GOLD)

    # Team
    add_picture(slide, ASSETS / "icon_users.png", 4.5, footer_top + 0.8,
                width_cm=3.2, height_cm=3.2)
    add_text(slide, 4.5, footer_top + 4.2, 26.0, 1.4,
             lines=[("Project Team", {"size_pt": 18, "bold": True,
                                       "color_hex": C_GOLD})], align="l")
    team_lines = [(name, {"size_pt": 14, "color_hex": C_WHITE,
                            "bold": True}) for name in TEAM]
    add_text(slide, 4.5, footer_top + 5.6, 26.0, 4.6,
             lines=team_lines, line_spacing=1.3, align="l")

    # Supervisor (middle)
    mid_x = (POSTER_W_CM - 9.0) / 2 + 0.5
    add_picture(slide, ASSETS / "icon_supervisor.png", mid_x, footer_top + 0.8,
                width_cm=3.2, height_cm=3.2)
    add_text(slide, mid_x, footer_top + 4.2, 24.0, 1.4,
             lines=[("Under Supervision",
                      {"size_pt": 18, "bold": True, "color_hex": C_GOLD})],
             align="l")
    add_text(slide, mid_x, footer_top + 5.6, 24.0, 1.4,
             lines=[(SUPERVISOR, {"size_pt": 17, "color_hex": C_WHITE,
                                    "bold": True})], align="l")
    add_text(slide, mid_x, footer_top + 7.1, 24.0, 1.4,
             lines=[("Project Supervisor", {"size_pt": 13,
                                              "color_hex": C_MUTED})],
             align="l")

    # College (right)
    col_x = POSTER_W_CM - 4.5 - 22.0
    add_picture(slide, ASSETS / "icon_college.png", col_x, footer_top + 0.8,
                width_cm=3.2, height_cm=3.2)
    add_text(slide, col_x, footer_top + 4.2, 22.0, 1.4,
             lines=[("Affiliation", {"size_pt": 18, "bold": True,
                                       "color_hex": C_GOLD})], align="l")
    add_text(slide, col_x, footer_top + 5.6, 22.0, 4.6,
             lines=[
                 ("College of Computing", {"size_pt": 14, "bold": True,
                                              "color_hex": C_WHITE}),
                 ("Department of Cybersecurity", {"size_pt": 14,
                                                     "color_hex": C_WHITE}),
                 ("Umm Al-Qura University — 2026", {"size_pt": 14,
                                                       "color_hex": C_WHITE,
                                                       "bold": True}),
             ], line_spacing=1.3, align="l")

    # ---------------------------------------------------------------- save
    out = Path(r"C:\Users\mohan\Desktop") / "CyberSatDetect_Poster_EN_Pro.pptx"
    prs.save(out)
    print("Saved:", out)
    return out


if __name__ == "__main__":
    build()
