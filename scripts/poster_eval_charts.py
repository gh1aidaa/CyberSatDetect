"""Generate poster-friendly evaluation charts using the REAL evaluation data
from ``thesis_official_evaluation_figures/`` and swap them into the Evaluation
section of ``3.pptx`` (replacing Picture 36 + Picture 196 + Picture 197).

The three charts produced are:

1.  ``poster_per_attack_bars.png``   -- big chart on the right side.
    Per-attack F1 / Recall / Balanced-Accuracy bars at the best-F1 threshold,
    from ``per_attack_full_metrics_4attacks.csv``.

2.  ``poster_confusion_matrix.png``  -- top small chart.
    Window-level confusion matrix at the best-F1 threshold
    (TP=54663 / TN=348449 / FP=5785 / FN=49), from
    ``confusion_matrix_summary_4attacks.csv``.

3.  ``poster_roc_pr_compact.png``    -- bottom small chart.
    Compact ROC + PR side-by-side panels reconstructed from the operating
    points in ``overall_threshold_metrics_4attacks.csv`` (AUC values come
    from ``evaluation_summary_4attacks.json``: ROC-AUC = 0.996, PR-AUC =
    0.956).

All charts use the poster's purple palette so they integrate cleanly with
the rest of the page without crowding.
"""
from __future__ import annotations

import csv
import json
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.util import Emu


# ============================================================================
# Paths
# ============================================================================
ROOT = Path(__file__).resolve().parents[1]
DATA_DIR = ROOT / "thesis_official_evaluation_figures"
OUT_DIR = ROOT / "scripts" / "poster_eval_assets"
OUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_PATH = Path(r"C:\Users\mohan\Desktop\3.pptx")


# ============================================================================
# Poster palette  (purple/navy with subtle gold accent)
# ============================================================================
C_PURPLE      = "#8B5BC4"
C_PURPLE_DEEP = "#4A2674"
C_PURPLE_MID  = "#6B3A9C"
C_PURPLE_SOFT = "#B591D9"
C_BG_TINT     = "#F6F2FB"
C_GOLD        = "#D49A1F"
C_GOLD_SOFT   = "#F2D89A"
C_GRID        = "#D5CCE0"
C_TEXT        = "#1F1830"
C_TEXT_MID    = "#574B66"
C_WHITE       = "#FFFFFF"


_BASE_RC = {
    "font.family": "DejaVu Sans",
    "axes.facecolor": C_WHITE,
    "figure.facecolor": C_WHITE,
    "axes.edgecolor": C_PURPLE_DEEP,
    "axes.labelcolor": C_TEXT,
    "axes.titlecolor": C_PURPLE_DEEP,
    "xtick.color": C_TEXT,
    "ytick.color": C_TEXT,
    "axes.spines.right": False,
    "axes.spines.top": False,
    "axes.grid": True,
    "grid.color": C_GRID,
    "grid.alpha": 0.85,
    "grid.linewidth": 0.7,
}


def _save(fig, path: Path, tight: bool = True) -> None:
    if tight:
        fig.savefig(path, dpi=220, facecolor=C_WHITE, bbox_inches="tight",
                    pad_inches=0.18)
    else:
        fig.savefig(path, dpi=220, facecolor=C_WHITE)
    plt.close(fig)


# ============================================================================
# Load real evaluation data
# ============================================================================
def _read_per_attack():
    """Return {attack: {f1, recall, balanced_accuracy, far, ...}} at best-F1."""
    out = {}
    with open(DATA_DIR / "per_attack_full_metrics_4attacks.csv",
              newline="", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            if row["threshold_name"] != "best_f1":
                continue
            out[row["attack_type"]] = {
                "f1": float(row["f1"]),
                "recall": float(row["recall"]),
                "balanced_accuracy": float(row["balanced_accuracy"]),
                "far": float(row["far"]),
                "precision": float(row["precision"]),
            }
    return out


def _read_confusion_best_f1():
    with open(DATA_DIR / "confusion_matrix_summary_4attacks.csv",
              newline="", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            if row["threshold_name"] == "best_f1":
                return {
                    "TP": int(row["TP"]), "TN": int(row["TN"]),
                    "FP": int(row["FP"]), "FN": int(row["FN"]),
                }
    raise RuntimeError("best_f1 row not found in confusion matrix CSV")


def _read_threshold_operating_points():
    """Read all operating points (threshold name + FAR + Recall + Precision)."""
    rows = []
    with open(DATA_DIR / "overall_threshold_metrics_4attacks.csv",
              newline="", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            rows.append({
                "name": row["threshold_name"],
                "threshold": float(row["threshold"]),
                "f1": float(row["f1"]),
                "precision": float(row["precision"]),
                "recall": float(row["recall"]),
                "far": float(row["far"]),
            })
    return rows


def _read_summary():
    with open(DATA_DIR / "evaluation_summary_4attacks.json",
              encoding="utf-8") as f:
        return json.load(f)


# ============================================================================
# Chart 1 -- Per-attack performance bars (big chart, right side)
# ============================================================================
def make_per_attack_bars(out_path: Path) -> None:
    plt.rcParams.update(_BASE_RC)
    per = _read_per_attack()
    attacks = ["drift", "freeze", "noise", "spike"]
    labels  = ["Drift", "Freeze", "Noise", "Spike"]
    f1   = [per[a]["f1"] for a in attacks]
    rec  = [per[a]["recall"] for a in attacks]
    bal  = [per[a]["balanced_accuracy"] for a in attacks]

    fig, ax = plt.subplots(figsize=(11.55, 5.0), dpi=220)
    x = np.arange(len(labels))
    w = 0.26

    b1 = ax.bar(x - w, f1, w, label="F1-Score",
                color=C_PURPLE, edgecolor=C_PURPLE_DEEP, linewidth=0.8)
    b2 = ax.bar(x,     rec, w, label="Recall",
                color=C_GOLD, edgecolor=C_PURPLE_DEEP, linewidth=0.8)
    b3 = ax.bar(x + w, bal, w, label="Balanced Acc.",
                color=C_PURPLE_SOFT, edgecolor=C_PURPLE_DEEP, linewidth=0.8)

    for bars in (b1, b2, b3):
        for r in bars:
            v = r.get_height()
            ax.text(r.get_x() + r.get_width() / 2, v + 0.012,
                    f"{v * 100:.1f}%", ha="center", va="bottom",
                    color=C_PURPLE_DEEP, fontsize=11, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=13, fontweight="bold",
                       color=C_PURPLE_DEEP)
    ax.set_yticks([0.5, 0.6, 0.7, 0.8, 0.9, 1.0])
    ax.set_yticklabels([f"{int(v * 100)}%" for v in [0.5, 0.6, 0.7, 0.8, 0.9, 1.0]],
                       fontsize=11)
    ax.set_ylim(0.45, 1.13)
    ax.set_ylabel("Score", fontsize=12, fontweight="bold",
                  color=C_PURPLE_DEEP)
    ax.set_title("Per-Attack Performance  (Best-F1 threshold = 0.0509)",
                 fontsize=14, fontweight="bold", color=C_PURPLE_DEEP, pad=12)

    leg = ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.13), ncol=3,
                    frameon=False, fontsize=12)
    for t in leg.get_texts():
        t.set_color(C_PURPLE_DEEP)
    ax.grid(axis="y", color=C_GRID, linewidth=0.7)
    for s in ax.spines.values():
        s.set_color(C_PURPLE_DEEP)
        s.set_alpha(0.4)

    plt.tight_layout()
    _save(fig, out_path)


# ============================================================================
# Chart 2 -- Confusion matrix at best-F1 threshold (compact heatmap)
# ============================================================================
def make_confusion_matrix(out_path: Path) -> None:
    plt.rcParams.update(_BASE_RC)
    cm = _read_confusion_best_f1()
    matrix = np.array([
        [cm["TN"], cm["FP"]],
        [cm["FN"], cm["TP"]],
    ], dtype=float)
    total = matrix.sum()
    pct = matrix / total * 100.0

    fig, ax = plt.subplots(figsize=(8.5, 4.9), dpi=220)
    # Use a single-purple gradient with explicit min/max alpha so all 4 cells
    # remain readable even when one count is enormous (TN >> others).
    cmap = matplotlib.colors.LinearSegmentedColormap.from_list(
        "purple_pos", [C_BG_TINT, C_PURPLE_SOFT, C_PURPLE, C_PURPLE_DEEP]
    )
    norm = matplotlib.colors.PowerNorm(gamma=0.35, vmin=0, vmax=matrix.max())
    ax.imshow(matrix, cmap=cmap, norm=norm, aspect="auto")

    for (i, j), v in np.ndenumerate(matrix):
        is_dark = norm(v) > 0.45
        txt_color = C_WHITE if is_dark else C_PURPLE_DEEP
        cell_label = ["TN", "FP", "FN", "TP"][i * 2 + j]
        ax.text(j, i - 0.18, cell_label,
                ha="center", va="center", color=txt_color,
                fontsize=14, fontweight="bold")
        ax.text(j, i + 0.05, f"{int(v):,}",
                ha="center", va="center", color=txt_color,
                fontsize=15, fontweight="bold")
        ax.text(j, i + 0.28, f"({pct[i, j]:.2f}%)",
                ha="center", va="center", color=txt_color,
                fontsize=11)

    ax.set_xticks([0, 1])
    ax.set_xticklabels(["Pred: Normal", "Pred: Attack"], fontsize=12,
                       fontweight="bold", color=C_PURPLE_DEEP)
    ax.set_yticks([0, 1])
    ax.set_yticklabels(["True: Normal", "True: Attack"], fontsize=12,
                       fontweight="bold", color=C_PURPLE_DEEP)
    ax.set_xticks(np.arange(-.5, 2, 1), minor=True)
    ax.set_yticks(np.arange(-.5, 2, 1), minor=True)
    ax.grid(which="minor", color=C_PURPLE_DEEP, linestyle="-",
            linewidth=1.2, alpha=0.35)
    ax.tick_params(which="minor", length=0)
    ax.set_title("Confusion Matrix  (Best-F1 threshold)",
                 fontsize=13, fontweight="bold", color=C_PURPLE_DEEP, pad=10)
    for s in ax.spines.values():
        s.set_color(C_PURPLE_DEEP)
        s.set_alpha(0.4)
    plt.tight_layout()
    _save(fig, out_path)


# ============================================================================
# Chart 3 -- Compact ROC + PR (using real operating points)
# ============================================================================
def _interp_curve(x_pts, y_pts):
    """Return a smoothly interpolated curve through (x, y) operating points."""
    order = np.argsort(x_pts)
    xs = np.asarray(x_pts)[order]
    ys = np.asarray(y_pts)[order]
    # De-duplicate same x to avoid interpolation glitches.
    xs_u, idx = np.unique(xs, return_index=True)
    ys_u = ys[idx]
    xs_dense = np.linspace(0, 1, 300)
    ys_dense = np.interp(xs_dense, xs_u, ys_u)
    return xs_dense, ys_dense


def make_roc_pr_compact(out_path: Path) -> None:
    plt.rcParams.update(_BASE_RC)
    rows = _read_threshold_operating_points()
    summary = _read_summary()

    # ROC points (FPR, TPR=Recall). FPR == FAR in this project.
    fpr_pts = [0.0] + [r["far"] for r in rows] + [1.0]
    tpr_pts = [0.0] + [r["recall"] for r in rows] + [1.0]

    # PR points (Recall, Precision). Anchor with (0, 1) on the left and
    # the prevalence baseline on the right so the curve drops gracefully.
    rec_pts = [0.0] + [r["recall"] for r in rows] + [1.0]
    prec_pts = [1.0] + [r["precision"] for r in rows] + [0.13]

    best = next(r for r in rows if r["name"] == "best_f1")

    fig, axes = plt.subplots(1, 2, figsize=(11.4, 5.32), dpi=220)
    fig.subplots_adjust(left=0.085, right=0.985, top=0.86, bottom=0.16,
                        wspace=0.32)

    # --- ROC (zoomed on the operationally relevant FAR <= 10% region) -------
    ax = axes[0]
    xs, ys = _interp_curve(fpr_pts, tpr_pts)
    ax.plot([0, 1], [0, 1], color=C_GRID, linestyle="--", linewidth=1.0,
            label="Random", zorder=1)
    ax.fill_between(xs, 0, ys, color=C_PURPLE, alpha=0.10, zorder=2)
    ax.plot(xs, ys, color=C_PURPLE, linewidth=2.6,
            label="CyberSatDetect", zorder=3)
    ax.scatter([best["far"]], [best["recall"]],
               color=C_GOLD, s=120, zorder=5, edgecolor=C_PURPLE_DEEP,
               linewidth=1.6, label="Best-F1")
    # Annotation kept INSIDE the visible window (xlim = 0..0.10).
    ax.annotate(
        f"Best-F1\nFAR = {best['far'] * 100:.2f}%\nRecall = {best['recall'] * 100:.2f}%",
        xy=(best["far"], best["recall"]),
        xytext=(0.045, 0.45),
        fontsize=10, color=C_PURPLE_DEEP, fontweight="bold",
        ha="left", va="center",
        arrowprops=dict(arrowstyle="-", color=C_PURPLE_DEEP, lw=0.9),
        zorder=6,
    )
    ax.set_xlim(0, 0.10)
    ax.set_ylim(0, 1.03)
    ax.set_xlabel("False Positive Rate", fontsize=11.5, fontweight="bold",
                  color=C_PURPLE_DEEP)
    ax.set_ylabel("True Positive Rate", fontsize=11.5, fontweight="bold",
                  color=C_PURPLE_DEEP)
    ax.set_title(f"ROC  -  AUC = {summary['curves']['roc_auc']:.3f}",
                 fontsize=12.5, fontweight="bold", color=C_PURPLE_DEEP, pad=8)
    ax.legend(loc="lower right", frameon=False, fontsize=9.5)
    for s in ax.spines.values():
        s.set_color(C_PURPLE_DEEP); s.set_alpha(0.4)

    # --- PR ------------------------------------------------------------------
    ax = axes[1]
    xs, ys = _interp_curve(rec_pts, prec_pts)
    ax.fill_between(xs, 0, ys, color=C_GOLD, alpha=0.12, zorder=2)
    ax.plot(xs, ys, color=C_GOLD, linewidth=2.6,
            label="CyberSatDetect", zorder=3)
    ax.scatter([best["recall"]], [best["precision"]],
               color=C_PURPLE, s=120, zorder=5, edgecolor=C_PURPLE_DEEP,
               linewidth=1.6, label="Best-F1")
    ax.annotate(
        f"Best-F1\nRecall = {best['recall'] * 100:.2f}%\nPrecision = {best['precision'] * 100:.2f}%",
        xy=(best["recall"], best["precision"]),
        xytext=(0.08, 0.34),
        fontsize=10, color=C_PURPLE_DEEP, fontweight="bold",
        ha="left", va="center",
        arrowprops=dict(arrowstyle="-", color=C_PURPLE_DEEP, lw=0.9),
        zorder=6,
    )
    ax.set_xlim(0, 1.0)
    ax.set_ylim(0, 1.03)
    ax.set_xlabel("Recall", fontsize=11.5, fontweight="bold",
                  color=C_PURPLE_DEEP)
    ax.set_ylabel("Precision", fontsize=11.5, fontweight="bold",
                  color=C_PURPLE_DEEP)
    ax.set_title(f"PR  -  AUC = {summary['curves']['pr_auc']:.3f}",
                 fontsize=12.5, fontweight="bold", color=C_PURPLE_DEEP, pad=8)
    ax.legend(loc="lower left", frameon=False, fontsize=9.5)
    for s in ax.spines.values():
        s.set_color(C_PURPLE_DEEP); s.set_alpha(0.4)

    _save(fig, out_path, tight=False)


# ============================================================================
# PPTX picture-replacement helper
# ============================================================================
def replace_picture_image(slide, picture_name: str, new_image_path: Path) -> bool:
    """Replace the underlying image of a Picture shape WHILE keeping its
    position and size. Returns True if the swap happened.
    """
    for shape in slide.shapes:
        if shape.name != picture_name:
            continue
        # `image` is a python-pptx Image part; we replace its blob.
        image_part = shape.image
        rId = shape._element.blip_rId
        slide_part = slide.part
        # Use the existing rId but point it at a new image part.
        from pptx.parts.image import ImagePart
        with open(new_image_path, "rb") as f:
            new_blob = f.read()
        new_image_part = ImagePart.new(slide_part.package, new_blob,
                                       Path(new_image_path).suffix.lstrip("."))
        slide_part.related_parts[rId] = new_image_part
        return True
    return False


def replace_picture_by_recreation(slide, picture_name: str,
                                  new_image_path: Path) -> bool:
    """Fallback: delete the old picture and add a new one at the SAME L/T/W/H.

    The picture's z-order changes (it moves to the end) but visually the box
    occupies exactly the same area.
    """
    target = None
    for shape in slide.shapes:
        if shape.name == picture_name:
            target = shape
            break
    if target is None:
        return False
    l, t, w, h = target.left, target.top, target.width, target.height
    sp = target._element
    sp.getparent().remove(sp)
    new_pic = slide.shapes.add_picture(str(new_image_path), l, t,
                                       width=w, height=h)
    new_pic.name = picture_name
    return True


# ============================================================================
# Driver
# ============================================================================
def main():
    print("Generating poster evaluation charts in", OUT_DIR)
    p_bars = OUT_DIR / "poster_per_attack_bars.png"
    p_cm = OUT_DIR / "poster_confusion_matrix.png"
    p_rocpr = OUT_DIR / "poster_roc_pr_compact.png"

    make_per_attack_bars(p_bars)
    make_confusion_matrix(p_cm)
    make_roc_pr_compact(p_rocpr)
    print("  -", p_bars.name)
    print("  -", p_cm.name)
    print("  -", p_rocpr.name)

    print(f"\nUpdating evaluation pictures in {PPTX_PATH} ...")
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides[0]

    swaps = [
        ("Picture 36",  p_bars),    # big right chart  (W=25.57, H=11.05 cm)
        ("Picture 196", p_cm),      # top small chart  (W=11.89, H=6.88 cm)
        ("Picture 197", p_rocpr),   # bot small chart  (W=9.04,  H=4.22 cm)
    ]
    for pic_name, img_path in swaps:
        ok = False
        try:
            ok = replace_picture_image(slide, pic_name, img_path)
        except Exception as e:
            print(f"  in-place swap failed for {pic_name}: {e}")
        if not ok:
            ok = replace_picture_by_recreation(slide, pic_name, img_path)
        print(f"  {pic_name}: {'OK' if ok else 'NOT FOUND'}  <-  {img_path.name}")

    prs.save(str(PPTX_PATH))
    print(f"Saved: {PPTX_PATH}")


if __name__ == "__main__":
    main()
