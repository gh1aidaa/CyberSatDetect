"""Inspect groups in the poster template."""
import glob, os
from pptx import Presentation
from pptx.util import Emu

matches = glob.glob(r"C:\Users\mohan\Desktop\*.pptx")
target = None
for m in matches:
    if "(3)" in m and "_" in os.path.basename(m):
        target = m
        break

prs = Presentation(target)
slide = prs.slides[0]


def walk(shape, depth=0):
    pad = "  " * depth
    try:
        left = Emu(shape.left or 0).cm
        top = Emu(shape.top or 0).cm
        w = Emu(shape.width or 0).cm
        h = Emu(shape.height or 0).cm
    except Exception:
        left = top = w = h = 0
    print(f"{pad}{shape.shape_type} name={shape.name!r} pos=({left:.2f},{top:.2f}) size=({w:.2f}x{h:.2f})")
    if shape.has_text_frame:
        for pi, para in enumerate(shape.text_frame.paragraphs):
            for ri, run in enumerate(para.runs):
                t = run.text.replace("\n", "\\n")
                fs = run.font.size
                print(f"{pad}  txt p{pi}r{ri}: {t!r} font={run.font.name} size={fs}")
    if hasattr(shape, "shapes"):
        for sub in shape.shapes:
            walk(sub, depth + 1)


for i, shape in enumerate(slide.shapes):
    if shape.shape_type and "GROUP" in str(shape.shape_type):
        print(f"\n=== Group shape {i}: {shape.name!r} ===")
        walk(shape)
