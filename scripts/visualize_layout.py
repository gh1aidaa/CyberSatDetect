"""Render a schematic layout of the rearranged poster to verify positioning.

Draws every shape as a rectangle on an A0 portrait canvas (84.1 x 118.9 cm),
colored by category. Saves a PNG image for quick visual inspection.
"""
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC = r"C:\Users\mohan\Desktop\3.pptx"
OUT_PNG = "poster3_layout_after.png"

SLIDE_W = 84.1
SLIDE_H = 118.9


def e2c(emu):
    return emu / 360000.0


def color_for(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return ("#1976d2", 0.35)
    if shape.has_text_frame and shape.text_frame.text.strip():
        return ("#2e7d32", 0.25)
    return ("#9e9e9e", 0.2)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[0]

    fig, ax = plt.subplots(figsize=(8.41, 11.89))
    ax.set_xlim(0, SLIDE_W)
    ax.set_ylim(SLIDE_H, 0)  # flip Y so origin is top-left
    ax.set_aspect("equal")
    ax.set_xticks(range(0, int(SLIDE_W) + 1, 5))
    ax.set_yticks(range(0, int(SLIDE_H) + 1, 5))
    ax.grid(True, color="#eeeeee", linewidth=0.4)

    # Safety margin: 5 cm
    safe = mpatches.Rectangle(
        (5, 5), SLIDE_W - 10, SLIDE_H - 10,
        fill=False, edgecolor="red", linestyle="--", linewidth=1.2
    )
    ax.add_patch(safe)

    for shape in slide.shapes:
        if shape.left is None:
            continue
        l, t = e2c(shape.left), e2c(shape.top)
        w, h = e2c(shape.width), e2c(shape.height)
        color, alpha = color_for(shape)
        ax.add_patch(mpatches.Rectangle(
            (l, t), w, h, facecolor=color, alpha=alpha,
            edgecolor="black", linewidth=0.3
        ))
        # Label small text content
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip().split("\n")[0]
            if txt and len(txt) < 40 and w > 3 and h > 1:
                ax.text(
                    l + w / 2, t + h / 2, txt[:30],
                    ha="center", va="center", fontsize=4, color="black"
                )

    ax.set_title("Rearranged Poster Layout (A0 portrait)  -  red dashed = 5cm safety margin")
    plt.tight_layout()
    plt.savefig(OUT_PNG, dpi=160, bbox_inches="tight")
    print(f"Saved {OUT_PNG}")


if __name__ == "__main__":
    main()
