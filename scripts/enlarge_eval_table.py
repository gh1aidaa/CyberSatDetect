"""Re-layout the Performance Metrics table in the Evaluation section of
``3.pptx`` so it fills the empty horizontal space that opened up after the
small charts (Picture 196 / 197) were removed.

The new layout uses *4 rows x 2 columns* (=8 metrics) instead of 8 stacked
rows. Each row spans the full width between the eval-panel left margin and
the big per-attack chart on the right.

We touch ONLY:
  * Rectangle  171 / 174 / 177 / 180          (kept and resized into 4 wide rows)
  * Rectangle  183 / 186 / 189 / 192          (DELETED -- no longer needed)
  * TextBox    170                            (header: stretched + larger font)
  * TextBox    172..194 (label/value pairs)   (re-positioned + larger fonts)

Nothing else (Picture 36, TextBox 195 test caption, eval panel, conclusion
section, etc.) is modified.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


PPTX = Path(r"C:\Users\mohan\Desktop\3.pptx")
EMU_PER_CM = 360000


def c2e(cm: float) -> int:
    return int(round(cm * EMU_PER_CM))


# ---------------------------------------------------------------------------
# Geometry (cm)
# ---------------------------------------------------------------------------
# Panel header (Rounded Rectangle 164) bottom = 82.42
# Picture 36 left edge                       = 47.84
# TextBox 195 (test caption) top             = 92.02
HEADER_TOP      = 82.45      # panel-header bottom = 82.42 (3 mm gap)
HEADER_HEIGHT   = 2.10       # was 1.55
TABLE_LEFT      = 5.70
TABLE_WIDTH     = 41.80      # right edge L = 47.50 -> 0.34 cm gap before Picture 36
ROW_HEIGHT      = 2.55       # was 1.95
NUM_ROWS        = 4
FIRST_ROW_TOP   = HEADER_TOP + HEADER_HEIGHT
# Resulting bottom: 82.45 + 2.10 + 4 * 2.55 = 94.75 cm   (eval-panel bottom = 94.92)

# Half-row layout (left half  +  right half)
HALF_WIDTH      = TABLE_WIDTH / 2
INNER_PAD_L     = 1.00
INNER_PAD_R     = 0.80
LBL_W_FRAC      = 0.65        # of (HALF_WIDTH - paddings)

# Pre-computed cell coordinates
LEFT_HALF_X     = TABLE_LEFT
RIGHT_HALF_X    = TABLE_LEFT + HALF_WIDTH

LEFT_LBL_X      = LEFT_HALF_X + INNER_PAD_L
LEFT_LBL_W      = (HALF_WIDTH - INNER_PAD_L - INNER_PAD_R) * LBL_W_FRAC
LEFT_VAL_X      = LEFT_LBL_X + LEFT_LBL_W
LEFT_VAL_W      = HALF_WIDTH - INNER_PAD_L - INNER_PAD_R - LEFT_LBL_W

RIGHT_LBL_X     = RIGHT_HALF_X + INNER_PAD_L
RIGHT_LBL_W     = LEFT_LBL_W
RIGHT_VAL_X     = RIGHT_LBL_X + RIGHT_LBL_W
RIGHT_VAL_W     = LEFT_VAL_W

# Font sizes (pt) — MAXIMUM growth to fill the entire eval-panel
HEADER_FONT_PT  = 32      # was 17 -> 24 -> 32
LABEL_FONT_PT   = 30      # was 16 -> 22 -> 30
VALUE_FONT_PT   = 40      # was 19 -> 28 -> 40


# ---------------------------------------------------------------------------
# Mapping:  row index  ->  (left label TB, left value TB, right label TB, right value TB)
# Metric pairs were chosen to group related KPIs together.
# ---------------------------------------------------------------------------
ROW_BG_NAMES = ["Rectangle 171", "Rectangle 174", "Rectangle 177", "Rectangle 180"]
ROW_BG_TO_DELETE = ["Rectangle 183", "Rectangle 186", "Rectangle 189", "Rectangle 192"]
HEADER_NAME = "TextBox 170"

ROW_CONTENT = [
    # row 0
    ("TextBox 172", "TextBox 173",   # F1-Score / 0.9493
     "TextBox 181", "TextBox 182"),  # Recall   / 0.9991
    # row 1
    ("TextBox 175", "TextBox 176",   # Accuracy / 98.57%
     "TextBox 184", "TextBox 185"),  # Precision/ 0.9043
    # row 2
    ("TextBox 178", "TextBox 179",   # Balanced Accuracy / 99.14%
     "TextBox 187", "TextBox 188"),  # FAR              / 1.63%
    # row 3
    ("TextBox 190", "TextBox 191",   # ROC-AUC / 0.9958
     "TextBox 193", "TextBox 194"),  # PR-AUC  / 0.9561
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def place(shape, left_cm, top_cm, width_cm, height_cm):
    shape.left = c2e(left_cm)
    shape.top = c2e(top_cm)
    shape.width = c2e(width_cm)
    shape.height = c2e(height_cm)


def set_runs_font(shape, *, size_pt=None, align=None, anchor_middle=True):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if anchor_middle:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in tf.paragraphs:
        if align is not None:
            para.alignment = align
        for run in para.runs:
            if size_pt is not None and run.text.strip():
                run.font.size = Pt(size_pt)


def delete_shape(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)


# ---------------------------------------------------------------------------
def main():
    prs = Presentation(str(PPTX))
    slide = prs.slides[0]

    # 1) Header --------------------------------------------------------------
    hdr = find_shape(slide, HEADER_NAME)
    if hdr is None:
        raise RuntimeError("TextBox 170 not found")
    place(hdr, TABLE_LEFT, HEADER_TOP, TABLE_WIDTH, HEADER_HEIGHT)
    set_runs_font(hdr, size_pt=HEADER_FONT_PT, align=PP_ALIGN.LEFT)

    # 2) Four row backgrounds + their cells ----------------------------------
    for r, bg_name in enumerate(ROW_BG_NAMES):
        bg = find_shape(slide, bg_name)
        if bg is None:
            print(f"  ! {bg_name} not found")
            continue
        row_top = FIRST_ROW_TOP + r * ROW_HEIGHT
        place(bg, TABLE_LEFT, row_top, TABLE_WIDTH, ROW_HEIGHT)

        left_lbl_name, left_val_name, right_lbl_name, right_val_name = ROW_CONTENT[r]
        left_lbl  = find_shape(slide, left_lbl_name)
        left_val  = find_shape(slide, left_val_name)
        right_lbl = find_shape(slide, right_lbl_name)
        right_val = find_shape(slide, right_val_name)

        # Left half
        if left_lbl is not None:
            place(left_lbl, LEFT_LBL_X, row_top, LEFT_LBL_W, ROW_HEIGHT)
            set_runs_font(left_lbl, size_pt=LABEL_FONT_PT, align=PP_ALIGN.LEFT)
        if left_val is not None:
            place(left_val, LEFT_VAL_X, row_top, LEFT_VAL_W, ROW_HEIGHT)
            set_runs_font(left_val, size_pt=VALUE_FONT_PT, align=PP_ALIGN.RIGHT)

        # Right half
        if right_lbl is not None:
            place(right_lbl, RIGHT_LBL_X, row_top, RIGHT_LBL_W, ROW_HEIGHT)
            set_runs_font(right_lbl, size_pt=LABEL_FONT_PT, align=PP_ALIGN.LEFT)
        if right_val is not None:
            place(right_val, RIGHT_VAL_X, row_top, RIGHT_VAL_W, ROW_HEIGHT)
            set_runs_font(right_val, size_pt=VALUE_FONT_PT, align=PP_ALIGN.RIGHT)

    # 3) Delete the four redundant row backgrounds ---------------------------
    for name in ROW_BG_TO_DELETE:
        sh = find_shape(slide, name)
        if sh is not None:
            delete_shape(slide, sh)
            print(f"  deleted {name}")
        else:
            print(f"  ! {name} not found (already removed?)")

    # 4) Report final geometry ----------------------------------------------
    bottom = FIRST_ROW_TOP + NUM_ROWS * ROW_HEIGHT
    print(f"\nTable bounds: T = {HEADER_TOP:.2f}..{bottom:.2f} cm  "
          f"L = {TABLE_LEFT:.2f}..{TABLE_LEFT + TABLE_WIDTH:.2f} cm   "
          f"(W={TABLE_WIDTH:.2f}, H={bottom - HEADER_TOP:.2f})")
    print("Picture 36 left edge ~ 47.84 cm  (table right edge "
          f"= {TABLE_LEFT + TABLE_WIDTH:.2f} cm).")
    print("Test-set caption top ~ 92.02 cm  (table bottom "
          f"= {bottom:.2f} cm).")

    prs.save(str(PPTX))
    print(f"\nSaved: {PPTX}")


if __name__ == "__main__":
    main()
