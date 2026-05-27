"""Normalize the layout of 3.pptx so that:

* All content boxes (1-6, 8 Conclusion, 9 Future Work) share the SAME size
  (W = 36.45 cm, H = 18.03 cm), aligned in a 2-column grid.
* The Evaluation box (7) stays BIG as is: full width (~74.1 cm) and its
  current height (15.5 cm).
* All accent bars, number badges and right-edge icons start/end at identical
  offsets from their box edges so they line up across all boxes.
* The purple/gold accent colours follow a clear alternating pattern
  (odd numbers = purple, even numbers = gold).
* Orphan / duplicate FUTURE-WORK shapes (Picture 216, TextBox 217) left over
  from an earlier layout are removed so nothing visually overlaps.
* The poster stays A0 portrait (84.1 x 118.9 cm) and respects a 5 cm safety
  margin from every side. The bottom row now uses the previously empty space
  so the poster looks balanced top to bottom.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

SRC = r"C:\Users\mohan\Desktop\3.pptx"
DST = r"C:\Users\mohan\Desktop\3.pptx"


# ============================================================================
# Geometry constants (all in centimetres)
# ============================================================================
SLIDE_W = 84.1
SLIDE_H = 118.9
SAFE = 5.0  # safety margin

# Two-column grid
COL1_L = 5.0
COL2_L = 42.65
BOX_W = 36.45
BOX_H = 18.03
HEADER_STRIP_H = 3.0

# Vertical layout: header + 3 rows + Eval + 1 row, evenly spaced.
HEADER_T = 4.0
HEADER_H = 17.61
ROW1_T = 22.54
ROW2_T = 41.50
ROW3_T = 60.46
EVAL_T  = 79.42
EVAL_H  = 15.5
EVAL_L  = 5.0
EVAL_W  = 74.1
ROW4_T  = 95.85   # ends at 113.88 (within 5cm bottom margin = 113.9)

# Standard sub-shape offsets (taken from boxes 1-6 in the source pptx).
BAR_W = 0.76
BAR_H = 1.8
BAR_OFFSET_L = 0.76
BAR_OFFSET_T = 0.6

BADGE_W = 3.36
BADGE_H = 1.9
BADGE_OFFSET_L = 1.98
BADGE_OFFSET_T = 0.55

ICON_W = 2.2
ICON_H = 2.2
ICON_OFFSET_R = 1.05   # offset of icon right edge from box right edge
ICON_OFFSET_T = 0.4

TITLE_OFFSET_L = 6.1
TITLE_OFFSET_T = 0.25
TITLE_H = 2.2


# ============================================================================
# Colour palette  -  unified PURPLE accents across every box; the number
# badge itself is WHITE with a purple outline (a more elegant look that
# matches the screenshot the user shared).
# ============================================================================
PURPLE = RGBColor(0x8B, 0x5B, 0xC4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# All boxes share the same accent colour now (purple). Kept as a map so the
# pattern can be tweaked later if needed.
BOX_COLOR = {bid: PURPLE for bid in range(1, 10)}

# Box-3 sub-pill buttons (DETECT / COVER / REDUCE / ADAPT / DEPLOY): originally
# teal (#0099B2); the user wants them recoloured to purple to match.
BOX3_SUB_BUTTONS = [
    "Rounded Rectangle 61",
    "Rounded Rectangle 64",
    "Rounded Rectangle 67",
    "Rounded Rectangle 70",
    "Rounded Rectangle 73",
]


# ============================================================================
# Per-box configuration  -  maps each logical box to the shape names that
# belong to it in the source pptx.
# ============================================================================
# Each entry:
#   old: (L, T, W, H) of the bounding rectangle in the SOURCE pptx
#   new: (L, T, W, H) of the bounding rectangle in the TARGET layout
#   bg:        name of the main rounded rectangle (the box background)
#   header:    name of the header strip
#   accent:    name of the left vertical accent bar
#   badge:     name of the number badge background (rounded rectangle)
#   badge_text: name of the textbox sitting on top of the badge (the digit)
#                 -- may be None if the digit lives inside `badge`
#   icon:      name of the small right-edge icon picture
#   title:     name of the header-strip title textbox
#   members:   list of ALL shape names inside this box (used for shifting
#              every internal content shape by the box's positional delta)
# ----------------------------------------------------------------------------
BOXES = {
    1: dict(
        old=(5.00, 22.70, 36.45, 18.03),
        new=(COL1_L, ROW1_T, BOX_W, BOX_H),
        bg="Rounded Rectangle 24", header="Rounded Rectangle 25",
        accent="Rectangle 26", badge="Rounded Rectangle 27",
        badge_text="TextBox 28", icon="Picture 29", title="TextBox 30",
        members=[
            "Rounded Rectangle 24", "Rounded Rectangle 25", "Rectangle 26",
            "Rounded Rectangle 27", "TextBox 28", "Picture 29", "TextBox 30",
            "TextBox 31", "TextBox 39",
        ],
    ),
    2: dict(
        old=(5.00, 41.93, 36.45, 18.03),
        new=(COL1_L, ROW2_T, BOX_W, BOX_H),
        bg="Rounded Rectangle 41", header="Rounded Rectangle 42",
        accent="Rectangle 43", badge="Rounded Rectangle 44",
        badge_text="TextBox 45", icon="Picture 46", title="TextBox 47",
        members=[
            "Rounded Rectangle 41", "Rounded Rectangle 42", "Rectangle 43",
            "Rounded Rectangle 44", "TextBox 45", "Picture 46", "TextBox 47",
            "TextBox 48",
        ],
    ),
    3: dict(
        old=(5.00, 61.17, 36.45, 18.03),
        new=(COL1_L, ROW3_T, BOX_W, BOX_H),
        bg="Rounded Rectangle 53", header="Rounded Rectangle 54",
        accent="Rectangle 55", badge="Rounded Rectangle 56",
        badge_text="TextBox 57", icon="Picture 58", title="TextBox 59",
        members=[
            "Rounded Rectangle 53", "Rounded Rectangle 54", "Rectangle 55",
            "Rounded Rectangle 56", "TextBox 57", "Picture 58", "TextBox 59",
            "TextBox 60",
            "Rounded Rectangle 61", "TextBox 62", "TextBox 63",
            "Rounded Rectangle 64", "TextBox 65", "TextBox 66",
            "Rounded Rectangle 67", "TextBox 68", "TextBox 69",
            "Rounded Rectangle 70", "TextBox 71", "TextBox 72",
            "Rounded Rectangle 73", "TextBox 74", "TextBox 75",
        ],
    ),
    4: dict(
        old=(42.65, 22.70, 36.45, 18.03),
        new=(COL2_L, ROW1_T, BOX_W, BOX_H),
        bg="Rounded Rectangle 76", header="Rounded Rectangle 77",
        accent="Rectangle 78", badge="Rounded Rectangle 79",
        badge_text="TextBox 80", icon="Picture 81", title="TextBox 82",
        members=[
            "Rounded Rectangle 76", "Rounded Rectangle 77", "Rectangle 78",
            "Rounded Rectangle 79", "TextBox 80", "Picture 81", "TextBox 82",
            "TextBox 253", "Picture 252",
        ],
    ),
    5: dict(
        old=(42.65, 41.93, 36.45, 18.03),
        new=(COL2_L, ROW2_T, BOX_W, BOX_H),
        bg="Rounded Rectangle 84", header="Rounded Rectangle 85",
        accent="Rectangle 86", badge="Rounded Rectangle 87",
        badge_text="TextBox 88", icon="Picture 89", title="TextBox 90",
        members=[
            "Rounded Rectangle 84", "Rounded Rectangle 85", "Rectangle 86",
            "Rounded Rectangle 87", "TextBox 88", "Picture 89", "TextBox 90",
            "TextBox 91", "Picture 294",
        ],
    ),
    6: dict(
        # NOTE: in the source the bg has W=36.64 (a typo); we treat it as 36.45.
        old=(42.65, 61.17, 36.45, 18.03),
        new=(COL2_L, ROW3_T, BOX_W, BOX_H),
        bg="Rounded Rectangle 93", header="Rounded Rectangle 94",
        accent="Rectangle 95", badge="Rounded Rectangle 96",
        badge_text="TextBox 97", icon="Picture 98", title="TextBox 99",
        members=[
            "Rounded Rectangle 93", "Rounded Rectangle 94", "Rectangle 95",
            "Rounded Rectangle 96", "TextBox 97", "Picture 98", "TextBox 99",
            "TextBox 254", "TextBox 257", "TextBox 258", "TextBox 259",
            "TextBox 260", "TextBox 261",
            "Picture 275", "Picture 277", "Picture 279", "Picture 288",
            "Picture 290", "Picture 292",
        ],
    ),
    7: dict(  # Evaluation — BIG, full-width
        old=(5.00, 80.00, 74.10, 15.50),
        new=(EVAL_L, EVAL_T, EVAL_W, EVAL_H),
        bg="Rounded Rectangle 163", header="Rounded Rectangle 164",
        accent="Rectangle 165",
        # For the Eval badge we need to pick the FIRST 'Rounded Rectangle 166'
        # (there are two shapes with this name in the source).
        badge="Rounded Rectangle 166",
        badge_text="TextBox 167", icon="Picture 168", title="TextBox 169",
        members=[
            "Rounded Rectangle 163", "Rounded Rectangle 164", "Rectangle 165",
            "Rounded Rectangle 166_eval",
            "TextBox 167", "Picture 168", "TextBox 169",
            "TextBox 170",
            "Rectangle 171", "TextBox 172", "TextBox 173",
            "Rectangle 174", "TextBox 175", "TextBox 176",
            "Rectangle 177", "TextBox 178", "TextBox 179",
            "Rectangle 180", "TextBox 181", "TextBox 182",
            "Rectangle 183", "TextBox 184", "TextBox 185",
            "Rectangle 186", "TextBox 187", "TextBox 188",
            "Rectangle 189", "TextBox 190", "TextBox 191",
            "Rectangle 192", "TextBox 193", "TextBox 194",
            "TextBox 195",
            "Picture 196", "Picture 197",
            "Picture 36",
        ],
    ),
    8: dict(  # Conclusion
        old=(4.69, 96.50, 36.93, 12.60),
        new=(COL1_L, ROW4_T, BOX_W, BOX_H),
        bg="Rounded Rectangle 201", header="Rounded Rectangle 202",
        accent="Rectangle 203", badge="Rounded Rectangle 204",
        badge_text="TextBox 205", icon="Picture 206", title="TextBox 207",
        members=[
            "Rounded Rectangle 201", "Rounded Rectangle 202", "Rectangle 203",
            "Rounded Rectangle 204", "TextBox 205", "Picture 206", "TextBox 207",
            "Rounded Rectangle 208", "TextBox 209", "TextBox 210",
        ],
    ),
    9: dict(  # Future Work
        old=(42.78, 96.05, 36.64, 12.60),
        new=(COL2_L, ROW4_T, BOX_W, BOX_H),
        bg="Rounded Rectangle 211", header="Rounded Rectangle 212",
        accent="Rectangle 115",
        # Box 9 stores the badge digit inside the rounded rectangle itself,
        # so badge and badge_text refer to the same shape.
        badge="Rounded Rectangle 166_fw",
        badge_text=None,
        icon="Picture 116", title="TextBox 117",
        members=[
            "Rounded Rectangle 211", "Rounded Rectangle 212", "Rectangle 115",
            "Rounded Rectangle 166_fw",
            "Picture 116", "TextBox 117", "TextBox 118",
        ],
    ),
}

# Orphan / duplicate shapes from earlier layout iterations.
#   Picture 216, TextBox 217 -- duplicate "FUTURE WORK" title at the OLD
#                               Future-Work position (T=101.7); now overlaps
#                               the body of the new FW box.
#   Picture 51, Picture 100  -- two small duplicate icons sitting at L≈65 in
#                               the Future-Work area; Picture 116 (already at
#                               L=75.85) is the legitimate top-right icon for
#                               that box, so 51 and 100 are leftover dups.
ORPHAN_NAMES = ["Picture 216", "TextBox 217", "Picture 51", "Picture 100"]


# ============================================================================
# Helpers
# ============================================================================
def c2e(cm):
    return int(round(cm * 360000.0))


def e2c(emu):
    return emu / 360000.0


def set_geom(shape, l=None, t=None, w=None, h=None):
    if l is not None: shape.left = c2e(l)
    if t is not None: shape.top = c2e(t)
    if w is not None: shape.width = c2e(w)
    if h is not None: shape.height = c2e(h)


def shift_shape(shape, dx, dy):
    if shape.left is not None:
        shape.left = int(shape.left + dx * 360000)
    if shape.top is not None:
        shape.top = int(shape.top + dy * 360000)


def index_shapes(slide):
    """Return a dict name -> list of shapes (handles duplicate names)."""
    idx = {}
    for shape in slide.shapes:
        idx.setdefault(shape.name, []).append(shape)
    return idx


def pick_rr166(slide, which):
    """Disambiguate the two 'Rounded Rectangle 166' shapes.

    `which` is 'eval' (the small badge near the Evaluation header) or 'fw'
    (the duplicate near the Future Work header).
    """
    candidates = [s for s in slide.shapes if s.name == "Rounded Rectangle 166"]
    candidates.sort(key=lambda s: e2c(s.left))
    if which == "eval":
        return candidates[0]   # smaller L
    return candidates[1]


def get_box_shape(slide, shape_index, box_id, key):
    """Return the shape object referenced by BOXES[box_id][key]."""
    name = BOXES[box_id].get(key)
    if not name:
        return None
    if name == "Rounded Rectangle 166_eval":
        return pick_rr166(slide, "eval")
    if name == "Rounded Rectangle 166_fw":
        return pick_rr166(slide, "fw")
    lst = shape_index.get(name, [])
    return lst[0] if lst else None


def iter_box_members(slide, shape_index, box_id):
    """Yield every shape that belongs to the given box (handles duplicate names)."""
    for member_name in BOXES[box_id]["members"]:
        if member_name == "Rounded Rectangle 166_eval":
            yield pick_rr166(slide, "eval")
        elif member_name == "Rounded Rectangle 166_fw":
            yield pick_rr166(slide, "fw")
        else:
            for s in shape_index.get(member_name, []):
                yield s


# ============================================================================
# Main transformation
# ============================================================================
def main():
    prs = Presentation(SRC)
    slide = prs.slides[0]

    # --- 1. Delete orphan / duplicate shapes ----------------------------------
    shape_index = index_shapes(slide)
    for orphan in ORPHAN_NAMES:
        for s in shape_index.get(orphan, []):
            sp = s._element
            sp.getparent().remove(sp)
        if shape_index.get(orphan):
            print(f"Deleted orphan: {orphan}")

    # Re-index after deletions.
    shape_index = index_shapes(slide)

    # --- 2. Move every box's members to their new positions. -----------------
    # For boxes whose height stays the same (1-6, 7) we do a simple shift.
    # For boxes whose height changes (8, 9: 12.6 -> 18.03) we ALSO scale the
    # vertical T-offset of every inner content shape so the contents spread
    # to fill the new taller box (avoids ~4-5 cm of empty whitespace at the
    # bottom). Shape heights are scaled only for the "body" content shapes;
    # the header strip / accent bar / badge / icon stay their normal size and
    # are repositioned by the standardisation step (#3) afterwards.
    STRUCT_SHAPES = {
        "Rounded Rectangle 201", "Rounded Rectangle 202", "Rectangle 203",
        "Rounded Rectangle 204", "TextBox 205", "Picture 206", "TextBox 207",
        "Rounded Rectangle 211", "Rounded Rectangle 212", "Rectangle 115",
        "Picture 116", "TextBox 117",
    }

    for box_id, cfg in BOXES.items():
        old_l, old_t, _, old_h = cfg["old"]
        new_l, new_t, _, new_h = cfg["new"]
        dx = new_l - old_l
        sy = new_h / old_h if old_h else 1.0
        for shape in iter_box_members(slide, shape_index, box_id):
            if shape is None:
                continue
            shape.left = int(shape.left + dx * 360000)
            # Vertical: keep structural shapes at original offsets (they will
            # be repositioned to standard offsets below). Body content shapes
            # have their offset scaled by sy so they spread across the new
            # box height.
            old_top_cm = e2c(shape.top)
            offset_t = old_top_cm - old_t
            if shape.name in STRUCT_SHAPES or sy == 1.0:
                new_top_cm = new_t + offset_t
            else:
                new_top_cm = new_t + offset_t * sy
                # For body textboxes/highlight bg, also scale height a bit so
                # the highlight bar widens vertically and matches the spread.
                if "Rounded Rectangle 208" in shape.name or "TextBox 209" in shape.name:
                    new_h_cm = e2c(shape.height) * sy
                    shape.height = c2e(new_h_cm)
            shape.top = c2e(new_top_cm)

    # --- 3. Standardise the box-defining shapes (bg, header, accent, badge,
    #        badge text, icon, title) so every box looks identical. ----------
    for box_id, cfg in BOXES.items():
        new_l, new_t, new_w, new_h = cfg["new"]

        # Background
        bg = get_box_shape(slide, shape_index, box_id, "bg")
        if bg is not None:
            set_geom(bg, l=new_l, t=new_t, w=new_w, h=new_h)

        # Header strip
        hdr = get_box_shape(slide, shape_index, box_id, "header")
        if hdr is not None:
            set_geom(hdr, l=new_l, t=new_t, w=new_w, h=HEADER_STRIP_H)

        # Accent bar
        accent = get_box_shape(slide, shape_index, box_id, "accent")
        if accent is not None:
            set_geom(
                accent,
                l=new_l + BAR_OFFSET_L,
                t=new_t + BAR_OFFSET_T,
                w=BAR_W, h=BAR_H,
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = BOX_COLOR[box_id]

        # Badge background -- white fill, purple outline.
        badge = get_box_shape(slide, shape_index, box_id, "badge")
        if badge is not None:
            set_geom(
                badge,
                l=new_l + BADGE_OFFSET_L,
                t=new_t + BADGE_OFFSET_T,
                w=BADGE_W, h=BADGE_H,
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = WHITE
            badge.line.color.rgb = BOX_COLOR[box_id]
            badge.line.width = Pt(1.5)

        # Badge text (may be the same as the badge shape for Box 9).
        # Re-colour the digit to purple so it contrasts against the white
        # badge background.
        badge_text = get_box_shape(slide, shape_index, box_id, "badge_text")
        if badge_text is not None and badge_text is not badge:
            set_geom(
                badge_text,
                l=new_l + BADGE_OFFSET_L,
                t=new_t + BADGE_OFFSET_T,
                w=BADGE_W, h=BADGE_H,
            )
        # Recolour the digit text (works for both the separate textbox and
        # the textbox embedded inside the badge shape for Box 9).
        for shape_with_digit in (badge_text, badge):
            if shape_with_digit is None or not shape_with_digit.has_text_frame:
                continue
            for para in shape_with_digit.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        run.font.color.rgb = BOX_COLOR[box_id]
                        break
                else:
                    continue
                break

        # Right-edge icon (keep its own W/H, just align it).
        icon = get_box_shape(slide, shape_index, box_id, "icon")
        if icon is not None:
            icon_w = e2c(icon.width) if icon.width else ICON_W
            icon_h = e2c(icon.height) if icon.height else ICON_H
            set_geom(
                icon,
                l=new_l + new_w - ICON_OFFSET_R - icon_w,
                t=new_t + ICON_OFFSET_T,
                w=icon_w, h=icon_h,
            )

        # Header title textbox
        title = get_box_shape(slide, shape_index, box_id, "title")
        if title is not None:
            title_w = new_w - (BADGE_OFFSET_L + BADGE_W + 0.6) - (ICON_OFFSET_R + ICON_W + 0.6)
            set_geom(
                title,
                l=new_l + BADGE_OFFSET_L + BADGE_W + 0.6,
                t=new_t + TITLE_OFFSET_T,
                w=title_w, h=TITLE_H,
            )

    # --- 4. Specific fixes for shapes that overflow their box. ---------------
    # TextBox 209 (Conclusion F1-stats text) is W=36.66 cm in the source,
    # which is wider than the box (W=36.45 cm) AND its original L=1.39 cm
    # makes the textbox extend past the left edge of the slide content area.
    # Re-anchor it so it matches the highlight background (Rounded Rectangle
    # 208) exactly and the text stays centred over the highlight bar.
    rr208 = shape_index.get("Rounded Rectangle 208", [None])[0]
    tb209 = shape_index.get("TextBox 209", [None])[0]
    if rr208 is not None and tb209 is not None:
        tb209.left = rr208.left
        tb209.top = rr208.top
        tb209.width = rr208.width
        tb209.height = rr208.height

    # --- 5. Recolour the Box-3 sub-pill buttons from teal to purple. ---------
    # In the source the five pill buttons DETECT / COVER / REDUCE / ADAPT /
    # DEPLOY are filled with #0099B2 (teal). The user wants them to match the
    # rest of the poster's purple accent.
    for btn_name in BOX3_SUB_BUTTONS:
        for shape in shape_index.get(btn_name, []):
            shape.fill.solid()
            shape.fill.fore_color.rgb = PURPLE

    prs.save(DST)
    print(f"Saved: {DST}")


if __name__ == "__main__":
    main()
