"""Rearrange CyberSatDetect_Poster_Academic.pptx:
 - Delete boxes 7, 8, 9 (right column).
 - Reorganize boxes 1-6 into a wider 2-column layout (H slightly enlarged).
 - Lift the upper sections to use empty space under the header.
 - Extend the Evaluation rectangle vertically so it visually contains its
   content (chart + KEY FINDINGS) and fully removes the 5.8 cm white gap
   that used to sit between the Evaluation box and the Conclusion row.
 - Move KEY FINDINGS section under the Performance-Metrics table so the
   bullets no longer overflow into the Conclusion area.
 - Preserve all content (text, images at original sizes); only reposition/resize containers.
 - Keep A0 portrait (84.1 x 118.9 cm) and respect a 5 cm safety margin on all sides.
"""
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC = r"C:\Users\mohan\Desktop\CyberSatDetect_Poster_Academic_BACKUP.pptx"
DST = r"C:\Users\mohan\Desktop\CyberSatDetect_Poster_Academic.pptx"

# Old layout (cm): each numbered content box position in the SOURCE pptx.
OLD_BOXES = {
    1: (5.0,  27.7,  23.9, 18.03),
    2: (5.0,  46.93, 23.9, 18.03),
    3: (5.0,  66.17, 23.9, 18.03),
    4: (30.1, 27.7,  23.9, 18.03),
    5: (30.1, 46.93, 23.9, 18.03),
    6: (30.1, 66.17, 23.9, 18.03),
    7: (55.2, 27.7,  23.9, 18.03),
    8: (55.2, 46.93, 23.9, 18.03),
    9: (55.2, 66.17, 23.9, 18.03),
}

# New layout: 2 columns wide (36.45 cm each) with 1.2 cm gap, fully inside 5cm safety margin.
# Left col @ L=5.0, right col @ L=42.65, both ending at 41.45 and 79.1 respectively.
# Box rows lifted UP and slightly TALLER (H 18.03 -> 18.5) to fill empty space.
NEW_W = 36.45
COL1_L = 5.0
COL2_L = 42.65
NEW_BOX_H = 18.5
ROW1_T = 22.5
ROW2_T = 42.0   # gap 1.0 cm from row 1 bottom (T=41.0)
ROW3_T = 61.5   # gap 1.0 cm from row 2 bottom (T=60.5)
NEW_BOXES = {
    1: (COL1_L, ROW1_T, NEW_W, NEW_BOX_H),
    2: (COL1_L, ROW2_T, NEW_W, NEW_BOX_H),
    3: (COL1_L, ROW3_T, NEW_W, NEW_BOX_H),
    4: (COL2_L, ROW1_T, NEW_W, NEW_BOX_H),
    5: (COL2_L, ROW2_T, NEW_W, NEW_BOX_H),
    6: (COL2_L, ROW3_T, NEW_W, NEW_BOX_H),
}

DELETE_BOX_IDS = {7, 8, 9}

# Evaluation box (Box 10): identified by shape name number range.
# In the source PPTX the Evaluation panel uses shape numbers 163..200
# (Rounded Rectangle 163/164, Rectangle 165, Rounded Rectangle 166,
#  TextBox 167-200, Picture 168/196/197/198, metric rectangles 171/174/...).
EVAL_NAME_NUM_RANGE = (163, 200)

# Evaluation panel target geometry.
# Source rect: L=5.0, T=85.0, W=74.1, H=15.5  -> bottom 100.5 (with content overflowing to ~112).
# Target rect: L=5.0, T=80.5, W=74.1, H=20.8  -> bottom 101.3 (matches Conclusion top exactly).
EVAL_SOURCE_T = 85.0
EVAL_NEW_T = 80.5
EVAL_NEW_H = 20.8
EVAL_SHIFT = EVAL_NEW_T - EVAL_SOURCE_T  # = -4.5 cm (move up)
EVAL_PANEL_NUM = 163  # Rounded Rectangle 163: the full Evaluation background panel
KF_TITLE_NUM = 199    # TextBox 199 - "KEY FINDINGS" heading
KF_BULLETS_NUM = 200  # TextBox 200 - KEY FINDINGS bullet list

# Relocate KEY FINDINGS to the left+middle area below the Performance Metrics
# table (where there was empty whitespace), so it no longer overflows into the
# Conclusion box.
KF_TITLE_NEW = (5.7, 95.5, 33.05, 1.2)     # L, T, W, H (cm)
KF_BULLETS_NEW = (5.7, 96.8, 33.05, 4.0)

SLIDE_W_CM = 84.1
SLIDE_H_CM = 118.9
SAFE_RIGHT = SLIDE_W_CM - 5.0  # 79.1
SAFE_LEFT = 5.0


def e2c(emu):
    return emu / 360000.0


def c2e(cm):
    return int(round(cm * 360000))


def shape_dims_cm(shape):
    return (e2c(shape.left), e2c(shape.top), e2c(shape.width), e2c(shape.height))


def find_box_id(shape):
    """Return which OLD box (if any) this shape belongs to.

    Uses center-point containment plus size filtering to skip slide-wide
    backgrounds (Rectangle 1, header decoration spanning all columns, etc.).
    """
    s_l, s_t, s_w, s_h = shape_dims_cm(shape)
    # Skip oversized shapes - they're slide-wide backgrounds/decorations.
    if s_w > 40 or s_h > 25:
        return None
    cx = s_l + s_w / 2
    cy = s_t + s_h / 2
    for bid, (bl, bt, bw, bh) in OLD_BOXES.items():
        if (bl - 0.25 <= cx <= bl + bw + 0.25 and
            bt - 0.25 <= cy <= bt + bh + 0.25):
            return bid
    return None


def transform_in_box(shape, bid):
    """Reposition a shape inside its (new) box.

    Rules:
      * Pictures keep their original width AND height (no stretching).
        - Header-icon pictures (top of box, near right) are anchored to the
          new right edge of the box.
        - Large content pictures (W > 15 cm) get horizontally centered in the
          new box.
        - Medium/small content pictures get their left position proportionally
          scaled, preserving alignment with their adjacent text labels.
      * Rectangles / text frames are scaled proportionally in the horizontal
        direction so panels widen with the box.
      * The vertical position offset (from the old box top) is scaled by the
        vertical box ratio so content fills the slightly taller new box.
      * The main box background rectangle (a shape whose W and H both match
        the OLD box width and height) is also stretched vertically to the
        NEW box height. All other shapes keep their original height.
      * Widths are clamped so nothing crosses the right safety margin (79.1).
    """
    bl_o, bt_o, bw_o, bh_o = OLD_BOXES[bid]
    bl_n, bt_n, bw_n, bh_n = NEW_BOXES[bid]
    sx = bw_n / bw_o
    sy = bh_n / bh_o

    s_l, s_t, s_w, s_h = shape_dims_cm(shape)
    is_pic = (shape.shape_type == MSO_SHAPE_TYPE.PICTURE)

    # ---- Horizontal positioning / sizing ----
    if is_pic:
        in_header = (s_t - bt_o) < 3.0
        near_right = ((bl_o + bw_o) - (s_l + s_w)) < 3.0
        if in_header and near_right:
            old_right_off = (bl_o + bw_o) - (s_l + s_w)
            new_l = (bl_n + bw_n) - old_right_off - s_w
            new_w = s_w
        elif s_w > 15.0:
            new_l = bl_n + bw_n / 2 - s_w / 2
            new_w = s_w
        else:
            offset_l = s_l - bl_o
            new_l = bl_n + offset_l * sx
            new_w = s_w
    else:
        offset_l = s_l - bl_o
        new_l = bl_n + offset_l * sx
        new_w = s_w * sx

    if new_l + new_w > SAFE_RIGHT + 0.05:
        new_w = max(0.4, SAFE_RIGHT - new_l)
    if new_l < SAFE_LEFT - 0.05:
        new_l = SAFE_LEFT

    # ---- Vertical positioning / sizing ----
    offset_t = s_t - bt_o
    # Scale T offset slightly so contents spread to fill the taller new box.
    new_t = bt_n + offset_t * sy

    # If this shape is the box's MAIN BACKGROUND rectangle (covers the entire
    # old box footprint), stretch its height to the new box height too.
    is_full_box_bg = (abs(s_w - bw_o) < 0.3) and (abs(s_h - bh_o) < 0.3)
    new_h = bh_n if is_full_box_bg else s_h

    shape.left = c2e(new_l)
    shape.width = c2e(new_w)
    shape.top = c2e(new_t)
    shape.height = c2e(new_h)


def _parse_shape_num(name):
    """Extract the trailing number from a shape name like 'Picture 198'."""
    if not name:
        return None
    m = re.findall(r"\d+", name)
    return int(m[-1]) if m else None


def is_eval_shape(shape):
    """Return True if the shape belongs to the Evaluation panel (Box 10).

    Detected by shape-name number being in EVAL_NAME_NUM_RANGE. This catches
    items that visually fall below the Evaluation rectangle (KEY FINDINGS
    title/text) which a pure spatial test would mis-classify as Conclusion.
    """
    n = _parse_shape_num(shape.name)
    if n is None:
        return False
    return EVAL_NAME_NUM_RANGE[0] <= n <= EVAL_NAME_NUM_RANGE[1]


def shift_shape_vertically(shape, delta_cm):
    """Move a shape up (negative delta) or down by delta_cm cm."""
    new_t = e2c(shape.top) + delta_cm
    shape.top = c2e(new_t)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[0]

    to_delete = []
    to_transform = []   # list of (shape, bid) for boxes 1-6
    to_handle_eval = [] # Evaluation-panel shapes

    for shape in slide.shapes:
        if shape.left is None or shape.top is None:
            continue
        # Boxes 1-9 detection by spatial position (in the OLD layout).
        bid = find_box_id(shape)
        if bid is not None:
            if bid in DELETE_BOX_IDS:
                to_delete.append(shape)
            elif bid in NEW_BOXES:
                to_transform.append((shape, bid))
            continue
        # Evaluation panel detection (by shape-name number range).
        if is_eval_shape(shape):
            to_handle_eval.append(shape)

    print(f"Deleting {len(to_delete)} shapes from boxes 7-9")
    print(f"Repositioning/resizing {len(to_transform)} shapes in boxes 1-6 "
          f"(H {18.03}->{NEW_BOX_H}cm, rows at T={ROW1_T}/{ROW2_T}/{ROW3_T})")
    print(f"Processing {len(to_handle_eval)} Evaluation-panel shapes "
          f"(shift {EVAL_SHIFT:+.2f}cm; extend panel H to {EVAL_NEW_H}cm; "
          f"relocate KEY FINDINGS under Performance Metrics)")

    for shape in to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

    for shape, bid in to_transform:
        transform_in_box(shape, bid)

    for shape in to_handle_eval:
        n = _parse_shape_num(shape.name)
        if n == KF_TITLE_NUM:
            # Relocate KEY FINDINGS title under the metrics table.
            l, t, w, h = KF_TITLE_NEW
            shape.left = c2e(l)
            shape.top = c2e(t)
            shape.width = c2e(w)
            shape.height = c2e(h)
        elif n == KF_BULLETS_NUM:
            # Relocate KEY FINDINGS bullets below the title.
            l, t, w, h = KF_BULLETS_NEW
            shape.left = c2e(l)
            shape.top = c2e(t)
            shape.width = c2e(w)
            shape.height = c2e(h)
        elif n == EVAL_PANEL_NUM:
            # Eval main background panel: shift up AND extend height
            # so it visually contains the chart and the relocated KEY FINDINGS.
            shift_shape_vertically(shape, EVAL_SHIFT)
            shape.height = c2e(EVAL_NEW_H)
        else:
            # All other Evaluation panel shapes: simple vertical shift.
            shift_shape_vertically(shape, EVAL_SHIFT)

    prs.save(DST)
    print(f"Saved: {DST}")


if __name__ == "__main__":
    main()
