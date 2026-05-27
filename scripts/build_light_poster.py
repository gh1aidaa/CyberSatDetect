# -*- coding: utf-8 -*-
"""Build the white academic CyberSatDetect poster (A0 portrait, IEEE-style)."""
from __future__ import annotations

import io
import os
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "scripts" / "poster_light"

# Academic palette  (purple primary)
C_WHITE       = "FFFFFF"
C_BG_TINT     = "F6F2FB"
C_BG_TINT2    = "EFE6F7"
C_NAVY        = "4A2674"      # primary deep purple (was navy)
C_NAVY_DEEP   = "351851"      # darker purple
C_BLUE        = "6B3A9C"      # medium purple (was blue)
C_ACCENT_BLUE = "8B5BC4"      # light purple (was accent blue)
C_CYAN        = "0099B2"
C_GOLD        = "D49A1F"
C_GREEN       = "1F9D55"
C_GRAY_DIV    = "D5CCE0"
C_GRAY_TINT   = "EAE2F1"
C_TEXT        = "1F1830"
C_TEXT_MID    = "574B66"

POSTER_W_CM = 84.1
POSTER_H_CM = 118.9
SAFE_MARGIN = 5.0


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def add_picture(slide, path, l, t, w=None, h=None,
                width_cm=None, height_cm=None):
    if w is None and width_cm is not None:
        w = width_cm
    if h is None and height_cm is not None:
        h = height_cm
    args = [str(path), Cm(l), Cm(t)]
    kw = {}
    if w is not None: kw["width"] = Cm(w)
    if h is not None: kw["height"] = Cm(h)
    return slide.shapes.add_picture(*args, **kw)


def add_rect(slide, l, t, w, h, *, fill_hex=None, outline_hex=None,
             outline_pt=0.0, rounded=False, corner=0.04):
    w = max(w, 0.1)
    h = max(h, 0.1)
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Cm(l), Cm(t), Cm(w), Cm(h))
    if rounded:
        shp.adjustments[0] = corner
    if fill_hex is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if outline_hex is None or outline_pt <= 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = RGBColor.from_string(outline_hex)
        shp.line.width = Pt(outline_pt)
    shp.shadow.inherit = False
    if shp.has_text_frame:
        for s in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(shp.text_frame, s, Cm(0))
    return shp


def add_text(slide, l, t, w, h, *, lines, color_hex=None, size_pt=14,
             bold=False, font_name="Calibri", align="l", v_align="t",
             line_spacing=None):
    w = max(w, 0.5)
    h = max(h, 0.5)
    tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for s in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, s, Cm(0))
    tf.auto_size = None
    tf.vertical_anchor = {"c": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM,
                           "t": MSO_ANCHOR.TOP}[v_align]
    al = {"l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT, "c": PP_ALIGN.CENTER,
          "j": PP_ALIGN.JUSTIFY}[align]

    def items_to_runs(item):
        if isinstance(item, str):
            return [(item, {})]
        if isinstance(item, tuple) and len(item) == 2 and isinstance(item[1], dict):
            return [item]
        return item  # already list of (text, style)

    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if line_spacing is not None:
            p.line_spacing = line_spacing
        runs = items_to_runs(item)
        for j, (chunk, style) in enumerate(runs):
            r = p.add_run()
            r.text = chunk
            r.font.name = style.get("font_name", font_name)
            r.font.size = Pt(style.get("size_pt", size_pt))
            r.font.bold = style.get("bold", bold)
            r.font.color.rgb = RGBColor.from_string(
                style.get("color_hex", color_hex or C_TEXT))
    return tb


def add_bullets(slide, l, t, w, h, lines, *, size_pt=15, color_hex=C_TEXT,
                line_spacing=1.25, bullet_color=C_BLUE, indent_cm=0.7):
    """Add bullet points with custom bullet markers."""
    for i, line in enumerate(lines):
        # Bullet diamond/square marker
        marker_y = t + i * (line_spacing * size_pt * 0.035) + 0.18
        # Use a textbox per line for accurate positioning
    # Simpler: one textbox with multiple paragraphs and a unicode bullet
    items = []
    for line in lines:
        items.append([("▪  ", {"color_hex": bullet_color, "bold": True,
                                "size_pt": size_pt}),
                       (line, {"color_hex": color_hex, "size_pt": size_pt})])
    return add_text(slide, l, t, w, h, lines=items, line_spacing=line_spacing)


# ---------------------------------------------------------------------------
# Section card (academic style)
# ---------------------------------------------------------------------------

def add_card(slide, l, t, w, h, *, number: str, title: str,
             icon_path: Path | None = None, header_color: str = C_NAVY,
             accent_color: str = C_ACCENT_BLUE) -> tuple[float, float]:
    # Border-only panel with thin outline
    add_rect(slide, l, t, w, h, fill_hex=C_WHITE,
             outline_hex=C_GRAY_DIV, outline_pt=1.0, rounded=True,
             corner=0.018)
    # Header bar
    hh = 3.0
    add_rect(slide, l, t, w, hh, fill_hex=header_color, rounded=True,
             corner=0.04)
    # Accent vertical stripe on left
    add_rect(slide, l + 0.5, t + 0.6, 0.5, hh - 1.2, fill_hex=accent_color)
    # Number badge
    if number:
        badge_w = 2.2
        add_rect(slide, l + 1.3, t + 0.55, badge_w, hh - 1.1,
                  fill_hex=C_WHITE, outline_hex=accent_color,
                  outline_pt=2.0, rounded=True, corner=0.18)
        add_text(slide, l + 1.3, t + 0.55, badge_w, hh - 1.1,
                  lines=[(number, {"size_pt": 24, "bold": True,
                                    "color_hex": header_color})],
                  align="c", v_align="c")
        title_x = l + 1.3 + badge_w + 0.7
    else:
        title_x = l + 1.4
    # Icon
    if icon_path is not None and icon_path.exists():
        icon_x = l + w - hh - 0.4
        add_picture(slide, icon_path, icon_x, t + 0.4,
                    width_cm=hh - 0.8, height_cm=hh - 0.8)
        max_title_w = (icon_x - title_x) - 0.4
    else:
        max_title_w = (l + w) - title_x - 0.4
    # Title text
    add_text(slide, title_x, t + 0.4, max_title_w, hh - 0.8,
             lines=[(title, {"size_pt": 22, "bold": True,
                              "color_hex": C_WHITE})],
             align="l", v_align="c")
    # Body region
    return (l + 0.7, t + hh + 0.5)


# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------
TITLE  = "CyberSatDetect"
SUB    = "AI-Based Satellite Telemetry Anomaly Detection System"
TAG    = "Hybrid LSTM-GRU Autoencoder  •  Self-Supervised  •  Continual Learning"

TEAM = [
    "Ghaidaa A. Algarni",
    "Alaa F. Alhazmi",
    "Bayan A. Alzahrani",
    "Ghaidaa A. Almasudi",
]
SUPERVISOR = "Dr. Ehad A. Aljarf"

INTRODUCTION = [
    "Satellites are critical global infrastructure for communication, navigation, Earth observation, and defense.",
    "Telemetry streams expose vital operational state and are an attractive target for cyber attackers.",
    "Mission failure, data loss, and degraded service are direct risks of undetected telemetry anomalies.",
]

PROBLEM = [
    "Rule-based monitoring misses subtle multi-channel attack patterns.",
    "Fixed thresholds produce high false-alarm rates and miss novel attacks.",
    "No safe way to keep a deployed model up-to-date with evolving operating regimes.",
    "Lack of an end-to-end AI defense layer for satellite telemetry.",
]

OBJECTIVES = [
    "Detect anomalies in satellite telemetry with an unsupervised hybrid model.",
    "Cover 4 attack types: Drift, Freeze, Noise, Spike.",
    "Reduce FAR using thresholds derived from normal-only data.",
    "Enable safe continual learning with human approval & rollback.",
    "Deliver dashboard, REST API, and full analysis pipeline.",
]

METHODOLOGY = [
    "Self-supervised: trained on normal telemetry only.",
    "Sliding windows  W = 100,  stride  S = 50  (50% overlap).",
    "Composite loss  L = W_recon·L_recon + W_pred·L_pred + W_grad·L_grad + W_sep·L_sep.",
    "Inference score  s = e_recon + e_pred + e_grad  (no separation term).",
    "Thresholds  p95 / p97 / p99 / p99.5 / p99.7 / 3σ  from normal-only data.",
    "Strict evaluation on attacked_v2 with 10% timestep-level rule.",
]

SYS_ARCH = [
    "FastAPI backend serving inference, evaluation, and continual-learning endpoints.",
    "React frontend for live monitoring, file analysis, and admin tools.",
    "Persistent storage for run artefacts and model registry (PENDING / APPROVED).",
    "Containerized deployment with Docker for reproducibility.",
]

HYBRID_MODEL = [
    "Encoder:  LSTM(64) → GRU(32)  bottleneck.",
    "Two heads:  Reconstruction  +  next-step Predictor.",
    "Separation loss enforces score margin between normal and pseudo-anomalies.",
    "Inference uses score = e_recon + e_pred + e_grad  only.",
]

TOOLS = [
    "Python  •  TensorFlow / Keras  •  NumPy  •  scikit-learn",
    "FastAPI  •  Uvicorn  •  Pydantic  •  REST APIs",
    "React  •  Recharts  •  TailwindCSS  •  Vite",
    "Docker  •  Git  •  GitHub Actions  •  PowerShell",
]

INTERFACES = [
    "Live Monitor  —  streams telemetry windows with real-time score & threshold overlay.",
    "File Analysis  —  CSV/NPY upload with per-channel anomaly reports & JSON summary.",
    "Reports Engine  —  exportable summaries per run.",
    "Continual-Learning Console  —  buffer → admin approve → fine-tune workflow.",
]

FEATURES = [
    "AI Anomaly Detection Pipeline (Hybrid AE + Predictor)",
    "Live Telemetry Monitoring & Alerting",
    "Multi-format File Analysis  (CSV / NPY)",
    "Continual Learning with Admin Approval",
    "Model Registry  +  one-click Rollback",
    "Token-secured REST API",
    "Space-Cybersecurity Dashboard",
]

METRICS_TABLE = [
    ("F1-Score",          "0.9493"),
    ("Accuracy",          "98.57%"),
    ("Balanced Accuracy", "99.14%"),
    ("Recall",            "0.9991"),
    ("Precision",         "0.9043"),
    ("FAR",               "1.63%"),
    ("ROC-AUC",           "0.9958"),
    ("PR-AUC",            "0.9561"),
]

CONCLUSION = [
    "CyberSatDetect achieves F1 = 0.949, ROC-AUC = 0.996, FAR ≈ 1.63% on a 4-attack benchmark.",
    "The hybrid LSTM-GRU model trained on normal-only data detects Drift, Freeze, Noise, and Spike attacks with high reliability.",
    "Continual learning with admin approval and model registry enables safe in-production updates without catastrophic forgetting.",
]

FUTURE = [
    "Multi-channel telemetry support",
    "Live operational evaluation on satellite missions",
    "Adaptive threshold policy that tracks distribution drift",
    "Compound attack scenarios  (Drift+Spike, Pattern Shift, Scale, Drop)",
    "Multi-task continual learning across diverse operating regimes",
]


# ---------------------------------------------------------------------------
# Build poster
# ---------------------------------------------------------------------------

def build() -> Path:
    prs = Presentation()
    prs.slide_width  = Cm(POSTER_W_CM)
    prs.slide_height = Cm(POSTER_H_CM)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 0) Pure white background
    add_rect(slide, 0, 0, POSTER_W_CM, POSTER_H_CM, fill_hex=C_WHITE)

    # Subtle decorative corner accents (top-left + bottom-right thin lines)
    for off in (0.7, 1.4, 2.1):
        add_rect(slide, 0, off, 14, 0.10,
                 fill_hex=C_ACCENT_BLUE if off == 0.7 else (
                     C_CYAN if off == 1.4 else C_GOLD))
    for off in (0.7, 1.4, 2.1):
        add_rect(slide, POSTER_W_CM - 14, POSTER_H_CM - off - 0.10, 14, 0.10,
                 fill_hex=C_ACCENT_BLUE if off == 0.7 else (
                     C_CYAN if off == 1.4 else C_GOLD))

    # ------------------------------------------------------------ HEADER
    H_TOP   = 4.0
    H_BOT   = 26.5
    H_H     = H_BOT - H_TOP

    # Header background bar (light tint with navy bottom rule)
    add_rect(slide, SAFE_MARGIN, H_TOP, POSTER_W_CM - 2 * SAFE_MARGIN,
             H_H, fill_hex=C_BG_TINT, outline_hex=C_GRAY_DIV,
             outline_pt=0.75)
    # Top accent strip
    add_rect(slide, SAFE_MARGIN, H_TOP, POSTER_W_CM - 2 * SAFE_MARGIN, 0.4,
             fill_hex=C_NAVY)
    # Bottom navy rule
    add_rect(slide, SAFE_MARGIN, H_BOT - 0.5,
             POSTER_W_CM - 2 * SAFE_MARGIN, 0.5, fill_hex=C_NAVY)
    # Three gold pips on bottom rule (decorative)
    pip_y = H_BOT - 0.95
    for i, cx in enumerate((POSTER_W_CM / 2 - 1.0, POSTER_W_CM / 2,
                             POSTER_W_CM / 2 + 1.0)):
        add_rect(slide, cx - 0.20, pip_y, 0.40, 0.40, fill_hex=C_GOLD,
                  rounded=True, corner=0.5)

    # UQU logo left and right
    logo_path = ASSETS / "uqu_logo.png"
    logo_h = H_H - 4.0
    if logo_path.exists():
        add_picture(slide, logo_path, SAFE_MARGIN + 1.0, H_TOP + 2.0,
                    height_cm=logo_h)
        add_picture(slide, logo_path,
                    POSTER_W_CM - SAFE_MARGIN - 1.0 - logo_h, H_TOP + 2.0,
                    height_cm=logo_h)

    # Title block (center between logos)
    title_left = SAFE_MARGIN + logo_h + 2.5
    title_w    = POSTER_W_CM - 2 * (SAFE_MARGIN + logo_h + 2.5)

    # University name label above title
    add_text(slide, title_left, H_TOP + 0.8, title_w, 1.4,
             lines=["UMM AL-QURA UNIVERSITY  •  COLLEGE OF COMPUTING  •  CYBERSECURITY DEPARTMENT"],
             color_hex=C_BLUE, size_pt=20, bold=True, align="c",
             font_name="Calibri")
    # Title
    add_text(slide, title_left, H_TOP + 2.6, title_w, 5.0,
             lines=[(TITLE, {"size_pt": 96, "bold": True,
                              "color_hex": C_NAVY})],
             align="c", v_align="t", line_spacing=1.0)
    # Subtitle
    add_text(slide, title_left, H_TOP + 9.3, title_w, 2.6,
             lines=[(SUB, {"size_pt": 32, "bold": True,
                            "color_hex": C_BLUE})],
             align="c", v_align="t", line_spacing=1.05)
    # Gold underline
    add_rect(slide, title_left + (title_w - 26) / 2, H_TOP + 12.3, 26, 0.22,
             fill_hex=C_GOLD)
    # Tagline
    add_text(slide, title_left, H_TOP + 12.7, title_w, 1.5,
             lines=[(TAG, {"size_pt": 22, "color_hex": C_TEXT_MID,
                            "bold": True})], align="c")

    # Team & supervisor (just above the navy rule)
    add_text(slide, title_left, H_TOP + 15.0, title_w, 1.4,
             lines=[("PROJECT TEAM", {"size_pt": 18, "bold": True,
                                       "color_hex": C_NAVY})], align="c")
    add_text(slide, title_left, H_TOP + 16.4, title_w, 1.6,
             lines=[("   •   ".join(TEAM),
                     {"size_pt": 24, "bold": True, "color_hex": C_TEXT})],
             align="c")
    add_text(slide, title_left, H_TOP + 18.4, title_w, 1.4,
             lines=[
                 ("SUPERVISOR:  ",
                  {"size_pt": 20, "bold": True, "color_hex": C_NAVY}),
                 (SUPERVISOR,
                  {"size_pt": 24, "bold": True, "color_hex": C_BLUE}),
             ], align="c")

    # ------------------------------------------------------------ BODY 3 cols
    BODY_TOP = H_BOT + 1.2
    eval_top = 85.0           # bottom band start
    body_bottom = eval_top - 0.8

    col_gap = 1.2
    col_w   = (POSTER_W_CM - 2 * SAFE_MARGIN - 2 * col_gap) / 3
    left_x   = SAFE_MARGIN
    center_x = SAFE_MARGIN + col_w + col_gap
    right_x  = SAFE_MARGIN + 2 * (col_w + col_gap)

    body_h = body_bottom - BODY_TOP   # ~58.3 cm
    row_gap = 1.2
    row_h   = (body_h - 2 * row_gap) / 3

    # --- LEFT COLUMN ---------------------------------------------------
    # 1. Introduction
    bx, by = add_card(slide, left_x, BODY_TOP, col_w, row_h,
                      number="1", title="INTRODUCTION",
                      icon_path=ASSETS / "icon_intro.png")
    add_bullets(slide, bx, by, col_w - 1.4, row_h - 5.5, INTRODUCTION,
                size_pt=21, line_spacing=1.5)
    # Key Stats strip at bottom of Intro card
    stats_y = BODY_TOP + row_h - 4.4
    stats = [("4", "ATTACK\nTYPES"), ("226K", "TEST\nWINDOWS"),
              ("0.949", "F1-SCORE\n(best-F1)")]
    chip_w = (col_w - 1.4 - 0.4 * 2) / 3
    for i, (val, lbl) in enumerate(stats):
        cx = left_x + 0.7 + i * (chip_w + 0.2)
        color = [C_BLUE, C_GOLD, C_CYAN][i]
        add_rect(slide, cx, stats_y, chip_w, 3.6,
                  fill_hex=C_BG_TINT, outline_hex=color, outline_pt=1.8,
                  rounded=True, corner=0.06)
        add_text(slide, cx, stats_y + 0.2, chip_w, 2.0,
                  lines=[(val, {"size_pt": 38, "bold": True,
                                  "color_hex": color})],
                  align="c", v_align="c")
        add_text(slide, cx, stats_y + 2.2, chip_w, 1.3,
                  lines=[(lbl, {"size_pt": 12, "bold": True,
                                  "color_hex": C_NAVY})],
                  align="c", v_align="c", line_spacing=1.0)

    # 2. Problem Statement
    bx, by = add_card(slide, left_x, BODY_TOP + (row_h + row_gap), col_w, row_h,
                      number="2", title="PROBLEM STATEMENT",
                      icon_path=ASSETS / "icon_problem.png",
                      header_color=C_NAVY, accent_color=C_GOLD)
    add_bullets(slide, bx, by, col_w - 1.4, row_h - 7.5, PROBLEM,
                size_pt=20, line_spacing=1.5, bullet_color=C_GOLD)
    # Before / After insight box
    ba_y = BODY_TOP + (row_h + row_gap) + row_h - 6.4
    add_rect(slide, left_x + 0.7, ba_y, col_w - 1.4, 5.6,
              fill_hex=C_BG_TINT, outline_hex=C_GOLD, outline_pt=1.5,
              rounded=True, corner=0.04)
    add_text(slide, left_x + 0.7, ba_y + 0.3, col_w - 1.4, 1.4,
             lines=[("WHY  AI-BASED  DETECTION?",
                     {"size_pt": 16, "bold": True, "color_hex": C_NAVY})],
             align="c", v_align="c")
    add_text(slide, left_x + 0.7, ba_y + 1.6, (col_w - 1.4) / 2 - 0.2, 4.0,
             lines=[
                 [("Rule-based:",
                    {"size_pt": 15, "bold": True, "color_hex": C_GOLD})],
                 [("✗ misses subtle patterns",
                    {"size_pt": 14, "color_hex": C_TEXT})],
                 [("✗ high false-alarm rate",
                    {"size_pt": 14, "color_hex": C_TEXT})],
                 [("✗ cannot adapt",
                    {"size_pt": 14, "color_hex": C_TEXT})],
             ], align="l", line_spacing=1.32)
    add_text(slide, left_x + 0.7 + (col_w - 1.4) / 2 + 0.2,
              ba_y + 1.6, (col_w - 1.4) / 2 - 0.2, 4.0,
             lines=[
                 [("Hybrid AI:",
                    {"size_pt": 15, "bold": True, "color_hex": C_GREEN})],
                 [("✓ learns from normal",
                    {"size_pt": 14, "color_hex": C_TEXT})],
                 [("✓ low FAR ≈ 1.6%",
                    {"size_pt": 14, "color_hex": C_TEXT})],
                 [("✓ continual learning",
                    {"size_pt": 14, "color_hex": C_TEXT})],
             ], align="l", line_spacing=1.32)

    # 3. Proposed Solution — solution pillars with tag badges
    bx, by = add_card(slide, left_x, BODY_TOP + 2 * (row_h + row_gap),
                      col_w, row_h,
                      number="3", title="PROPOSED  SOLUTION",
                      icon_path=ASSETS / "icon_solution.png",
                      header_color=C_NAVY, accent_color=C_CYAN)
    # Solution headline
    add_text(slide, bx, by, col_w - 1.4, 2.2,
             lines=[("A self-supervised Hybrid LSTM-GRU Autoencoder that learns from normal-only telemetry and flags anomalies in real time.",
                      {"size_pt": 16, "bold": True, "color_hex": C_NAVY})],
             align="l", v_align="t", line_spacing=1.3)
    obj_items = [
        ("DETECT",  "Unsupervised hybrid model for anomaly detection."),
        ("COVER",   "Four attack types: Drift, Freeze, Noise, Spike."),
        ("REDUCE",  "Lower FAR using normal-only thresholds."),
        ("ADAPT",   "Continual learning with admin approval."),
        ("DEPLOY",  "Dashboard, REST API, and offline analysis."),
    ]
    obj_y = by + 2.6
    obj_row_h = (row_h - 7.0) / len(obj_items)
    for tag, text in obj_items:
        add_rect(slide, bx, obj_y + 0.05, 3.0, obj_row_h - 0.2,
                  fill_hex=C_CYAN, rounded=True, corner=0.18)
        add_text(slide, bx, obj_y + 0.05, 3.0, obj_row_h - 0.2,
                  lines=[(tag, {"size_pt": 13, "bold": True,
                                  "color_hex": C_WHITE})],
                  align="c", v_align="c")
        add_text(slide, bx + 3.3, obj_y, col_w - 4.7, obj_row_h,
                  lines=[(text, {"size_pt": 16, "color_hex": C_TEXT})],
                  align="l", v_align="c")
        obj_y += obj_row_h

    # --- CENTER COLUMN -------------------------------------------------
    # 4. Methodology + AI Pipeline
    bx, by = add_card(slide, center_x, BODY_TOP, col_w, row_h,
                      number="4", title="METHODOLOGY  &  AI PIPELINE",
                      icon_path=ASSETS / "icon_methodology.png")
    method_text_h = row_h - 6.5  # leave room for pipeline diagram
    add_bullets(slide, bx, by, col_w - 1.4, method_text_h, METHODOLOGY,
                size_pt=16, line_spacing=1.35)

    # 5. System Architecture (bullets first, then diagram fills remaining)
    bx, by = add_card(slide, center_x, BODY_TOP + (row_h + row_gap),
                      col_w, row_h,
                      number="5", title="SYSTEM ARCHITECTURE",
                      icon_path=ASSETS / "icon_architecture.png")
    bullet_block_h = 5.6
    add_bullets(slide, bx, by, col_w - 1.4, bullet_block_h, SYS_ARCH,
                size_pt=16, line_spacing=1.32)
    arch_path = ASSETS / "architecture.png"
    if arch_path.exists():
        avail_h = (BODY_TOP + (row_h + row_gap) + row_h) - (by + bullet_block_h + 0.3) - 0.5
        arch_w = col_w - 1.0
        arch_h_by_w = arch_w * (7 / 11)
        if arch_h_by_w > avail_h:
            arch_h = avail_h
            arch_w_adj = arch_h * (11 / 7)
            arch_x = center_x + (col_w - arch_w_adj) / 2
            arch_w = arch_w_adj
        else:
            arch_h = arch_h_by_w
            arch_x = center_x + 0.5
        add_picture(slide, arch_path, arch_x,
                    by + bullet_block_h + 0.3, w=arch_w, h=arch_h)

    # 6. Hybrid LSTM-GRU Model (bullets first, diagram fills rest)
    bx, by = add_card(slide, center_x, BODY_TOP + 2 * (row_h + row_gap),
                      col_w, row_h,
                      number="6", title="HYBRID  LSTM-GRU  MODEL",
                      icon_path=ASSETS / "icon_brain.png")
    bullet_block_h = 5.0
    add_bullets(slide, bx, by, col_w - 1.4, bullet_block_h, HYBRID_MODEL,
                size_pt=16, line_spacing=1.32)
    hybrid_path = ASSETS / "hybrid_model.png"
    if hybrid_path.exists():
        avail_h = (BODY_TOP + 2 * (row_h + row_gap) + row_h) - (by + bullet_block_h + 0.3) - 0.5
        hy_w = col_w - 1.0
        hy_h_by_w = hy_w * (5.5 / 10)
        if hy_h_by_w > avail_h:
            hy_h = avail_h
            hy_w_adj = hy_h * (10 / 5.5)
            hy_x = center_x + (col_w - hy_w_adj) / 2
            hy_w = hy_w_adj
        else:
            hy_h = hy_h_by_w
            hy_x = center_x + 0.5
        add_picture(slide, hybrid_path, hy_x,
                    by + bullet_block_h + 0.3, w=hy_w, h=hy_h)

    # --- RIGHT COLUMN --------------------------------------------------
    # 7. Tools & Technologies
    bx, by = add_card(slide, right_x, BODY_TOP, col_w, row_h,
                      number="7", title="TOOLS  &  TECHNOLOGIES",
                      icon_path=ASSETS / "icon_tools.png",
                      header_color=C_NAVY, accent_color=C_CYAN)
    # Category headers + chip rows
    chip_groups = [
        ("AI / DATA",      C_BLUE,    [("Python",), ("TensorFlow",),
                                         ("NumPy",), ("scikit-learn",)]),
        ("BACKEND",        C_CYAN,    [("FastAPI",), ("Uvicorn",),
                                         ("REST API",), ("Pydantic",)]),
        ("FRONTEND",       C_GOLD,    [("React",), ("Recharts",),
                                         ("Tailwind",), ("Vite",)]),
        ("DEVOPS",         C_GREEN,   [("Docker",), ("Git",),
                                         ("GitHub Actions",), ("PowerShell",)]),
    ]
    chip_y = by
    chip_h = 1.4
    group_gap = 0.55
    for cat_name, color, row in chip_groups:
        add_text(slide, bx, chip_y, col_w - 1.4, 1.0,
                  lines=[(cat_name, {"size_pt": 14, "bold": True,
                                       "color_hex": color})],
                  align="l", v_align="c")
        chip_y += 1.0
        chip_x = bx
        avail_w = col_w - 1.4
        total_w = sum(0.36 * len(label[0]) + 1.4 for label in row) + 0.25 * (len(row) - 1)
        scale = min(1.0, avail_w / total_w)
        for (label,) in row:
            cw = (0.36 * len(label) + 1.4) * scale
            add_rect(slide, chip_x, chip_y, cw, chip_h,
                      fill_hex=C_WHITE, outline_hex=color, outline_pt=2.0,
                      rounded=True, corner=0.3)
            add_text(slide, chip_x, chip_y, cw, chip_h,
                      lines=[(label, {"size_pt": 14, "bold": True,
                                       "color_hex": color})],
                      align="c", v_align="c")
            chip_x += cw + 0.25 * scale
        chip_y += chip_h + group_gap

    # 8. Implementation / Interfaces
    bx, by = add_card(slide, right_x, BODY_TOP + (row_h + row_gap),
                      col_w, row_h,
                      number="8", title="IMPLEMENTATION  &  INTERFACES",
                      icon_path=ASSETS / "icon_screen.png",
                      header_color=C_NAVY, accent_color=C_ACCENT_BLUE)
    dash_path = ASSETS / "dashboard_mock.png"
    if dash_path.exists():
        dh_w = col_w - 1.0
        dh_h = dh_w * (5.5 / 11)
        add_picture(slide, dash_path, right_x + 0.5, by, w=dh_w, h=dh_h)
        remaining = (BODY_TOP + (row_h + row_gap) + row_h) - (by + dh_h + 0.5) - 0.4
        if remaining > 1.0:
            add_bullets(slide, bx, by + dh_h + 0.4, col_w - 1.4,
                         remaining, INTERFACES,
                         size_pt=15, line_spacing=1.32)
    else:
        add_bullets(slide, bx, by, col_w - 1.4, row_h - 3.0, INTERFACES,
                     size_pt=18)

    # 9. System Modules / Features
    bx, by = add_card(slide, right_x, BODY_TOP + 2 * (row_h + row_gap),
                      col_w, row_h,
                      number="9", title="SYSTEM  MODULES",
                      icon_path=ASSETS / "icon_rocket.png",
                      header_color=C_NAVY, accent_color=C_GOLD)
    items = [[(f"✓  ", {"size_pt": 24, "bold": True,
                         "color_hex": C_GREEN}),
                (f"{m}", {"size_pt": 19, "color_hex": C_TEXT})]
               for m in FEATURES]
    add_text(slide, bx, by, col_w - 1.4, row_h - 4.0,
              lines=items, line_spacing=1.55)

    # ----------------- Pipeline diagram inside center card 4 ---------------
    pipe_path = ASSETS / "pipeline.png"
    if pipe_path.exists():
        pipe_w = col_w - 1.0
        pipe_h = pipe_w * (2.8 / 11)
        add_picture(slide, pipe_path, center_x + 0.5,
                    BODY_TOP + row_h - pipe_h - 0.4, w=pipe_w, h=pipe_h)

    # ------------------------------------------------------------ EVALUATION
    EV_TOP = eval_top
    EV_BOT = 100.5
    EV_H   = EV_BOT - EV_TOP
    ev_w = POSTER_W_CM - 2 * SAFE_MARGIN
    bx, by = add_card(slide, SAFE_MARGIN, EV_TOP, ev_w, EV_H,
                       number="10",
                       title="EVALUATION  &  RESULTS",
                       icon_path=ASSETS / "icon_chart.png",
                       header_color=C_NAVY, accent_color=C_GOLD)

    inner_left = bx
    inner_top  = by
    inner_w    = ev_w - 1.4

    # Left sub-area: Metrics table  (28% wide)
    tbl_w = inner_w * 0.28
    tbl_title_h = 1.2
    add_text(slide, inner_left, inner_top, tbl_w, tbl_title_h,
              lines=[("Performance Metrics  (best-F1)",
                       {"size_pt": 17, "bold": True, "color_hex": C_NAVY})],
              align="l")
    # Notes block height reserved
    notes_h = 2.4
    tbl_h = EV_H - 4.0 - tbl_title_h - notes_h - 0.6
    row_h_tbl = tbl_h / len(METRICS_TABLE)
    tbl_y0 = inner_top + tbl_title_h + 0.2
    for i, (k, v) in enumerate(METRICS_TABLE):
        y = tbl_y0 + i * row_h_tbl
        bg = C_BG_TINT if i % 2 == 0 else C_WHITE
        add_rect(slide, inner_left, y, tbl_w, row_h_tbl,
                  fill_hex=bg, outline_hex=C_GRAY_DIV, outline_pt=0.7)
        add_text(slide, inner_left + 0.5, y, tbl_w * 0.55, row_h_tbl,
                  lines=[(k, {"size_pt": 16, "color_hex": C_TEXT_MID,
                              "bold": True})],
                  align="l", v_align="c")
        add_text(slide, inner_left + tbl_w * 0.55, y,
                  tbl_w * 0.40, row_h_tbl,
                  lines=[(v, {"size_pt": 19, "bold": True,
                                "color_hex": C_NAVY})],
                  align="r", v_align="c")
    # Notes under table
    notes_y = tbl_y0 + tbl_h + 0.4
    add_text(slide, inner_left, notes_y, tbl_w, notes_h,
              lines=[
                  [("Test set:", {"size_pt": 14, "bold": True,
                                    "color_hex": C_NAVY}),
                   (" 72 files  •  226,623 normal windows",
                      {"size_pt": 13, "color_hex": C_TEXT_MID})],
                  [("Attacks: ", {"size_pt": 14, "bold": True,
                                    "color_hex": C_NAVY}),
                   ("Drift, Freeze, Noise, Spike",
                      {"size_pt": 13, "color_hex": C_TEXT_MID})],
                  [("Rule: ", {"size_pt": 14, "bold": True,
                                  "color_hex": C_NAVY}),
                   ("strict 10% timestep-level",
                      {"size_pt": 13, "color_hex": C_TEXT_MID})],
              ], align="l", line_spacing=1.35)

    # Middle sub-area: per-attack bar chart  (cap height to fit)
    bar_x = inner_left + tbl_w + 0.8
    bar_w_ideal = inner_w * 0.36
    bar_ratio = 5.5 / 9.5
    # Reserve space for score distribution below
    sd_ratio = 4.2 / 9
    avail_total = EV_H - 4.0   # full inner height
    # Allocate ~ 65% to bar, 35% to score distribution (plus small gap)
    gap = 0.4
    bar_h = min(bar_w_ideal * bar_ratio, (avail_total - gap) * 0.62)
    bar_w = min(bar_w_ideal, bar_h / bar_ratio)
    bars_path = ASSETS / "per_attack_bars_light.png"
    if bars_path.exists():
        add_picture(slide, bars_path, bar_x, inner_top, w=bar_w, h=bar_h)

    # Below the bar chart: score distribution mini
    sd_path = ASSETS / "score_distribution.png"
    if sd_path.exists():
        sd_h = max(2.0, avail_total - bar_h - gap)
        sd_w = min(bar_w, sd_h * (9 / 4.2))
        sd_h = min(sd_h, sd_w * (4.2 / 9))
        add_picture(slide, sd_path, bar_x, inner_top + bar_h + gap,
                    w=sd_w, h=sd_h)

    # Right sub-area: ROC + PR
    roc_x = bar_x + bar_w + 0.8
    roc_w = inner_left + inner_w - roc_x
    roc_h = roc_w * (5.4 / 12.5)
    roc_path = ASSETS / "roc_pr_combined.png"
    if roc_path.exists():
        add_picture(slide, roc_path, roc_x, inner_top, w=roc_w, h=roc_h)
    # Key findings under ROC/PR
    add_text(slide, roc_x, inner_top + roc_h + 0.6, roc_w, 1.2,
              lines=[("KEY FINDINGS",
                       {"size_pt": 18, "bold": True, "color_hex": C_NAVY})],
              align="l")
    findings = [
        "Near-perfect separation: normal vs. attacked windows.",
        "ROC-AUC = 0.996 confirms robust threshold-free ranking.",
        "Best-F1 point: recall ≈ 0.999 at FAR = 1.63%.",
        "Continual learning preserves accuracy across regime drift.",
    ]
    add_bullets(slide, roc_x, inner_top + roc_h + 1.9, roc_w,
                 EV_H - (roc_h + 6.0), findings,
                 size_pt=16, line_spacing=1.38, bullet_color=C_GOLD)

    # ------------------------------------------------------------ BOTTOM ROW
    BR_TOP = EV_BOT + 0.8
    BR_BOT = POSTER_H_CM - SAFE_MARGIN
    BR_H   = BR_BOT - BR_TOP

    qr_w = 9.5
    bottom_gap = 1.0
    main_w = POSTER_W_CM - 2 * SAFE_MARGIN - qr_w - bottom_gap
    half_w = (main_w - bottom_gap) / 2

    # Conclusion card
    bx, by = add_card(slide, SAFE_MARGIN, BR_TOP, half_w, BR_H,
                       number="11", title="CONCLUSION",
                       icon_path=ASSETS / "icon_solution.png",
                       header_color=C_NAVY, accent_color=C_GREEN)
    # Big highlight metric strip
    high_h = 2.6
    add_rect(slide, bx, by, half_w - 1.4, high_h, fill_hex=C_BG_TINT,
              outline_hex=C_GREEN, outline_pt=1.6, rounded=True, corner=0.08)
    add_text(slide, bx, by, half_w - 1.4, high_h,
             lines=[[("F1 = 0.949   ",
                       {"size_pt": 24, "bold": True, "color_hex": C_NAVY}),
                      ("•   ROC-AUC = 0.996   ",
                       {"size_pt": 22, "bold": True, "color_hex": C_BLUE}),
                      ("•   FAR ≈ 1.63%",
                       {"size_pt": 22, "bold": True, "color_hex": C_GOLD})]],
             align="c", v_align="c")
    add_bullets(slide, bx, by + high_h + 0.3, half_w - 1.4,
                 BR_H - high_h - 4.0, CONCLUSION,
                 size_pt=16, line_spacing=1.38, bullet_color=C_GREEN)

    # Future Work card
    fw_x = SAFE_MARGIN + half_w + bottom_gap
    bx, by = add_card(slide, fw_x, BR_TOP, half_w, BR_H,
                       number="12", title="FUTURE WORK",
                       icon_path=ASSETS / "icon_rocket.png",
                       header_color=C_NAVY, accent_color=C_GOLD)
    add_bullets(slide, bx, by, half_w - 1.4, BR_H - 4.0, FUTURE,
                 size_pt=17, line_spacing=1.42, bullet_color=C_GOLD)

    # QR card
    qr_x = POSTER_W_CM - SAFE_MARGIN - qr_w
    add_rect(slide, qr_x, BR_TOP, qr_w, BR_H,
              fill_hex=C_WHITE, outline_hex=C_GRAY_DIV, outline_pt=1.0,
              rounded=True, corner=0.018)
    # Top band
    add_rect(slide, qr_x, BR_TOP, qr_w, 2.3, fill_hex=C_NAVY, rounded=True,
              corner=0.04)
    add_text(slide, qr_x, BR_TOP + 0.4, qr_w, 1.5,
              lines=[("SCAN FOR PROJECT INFO",
                       {"size_pt": 16, "bold": True, "color_hex": C_WHITE})],
              align="c", v_align="c")
    qr_path = ASSETS / "qr.png"
    if qr_path.exists():
        qr_box = min(qr_w - 1.6, BR_H - 5.0)
        add_picture(slide, qr_path, qr_x + (qr_w - qr_box) / 2,
                     BR_TOP + 2.7, w=qr_box, h=qr_box)
        add_text(slide, qr_x + 0.4, BR_TOP + 2.7 + qr_box + 0.2,
                  qr_w - 0.8, 1.5,
                  lines=[("CyberSatDetect", {"size_pt": 16, "bold": True,
                                              "color_hex": C_NAVY})],
                  align="c")

    # Save
    out = Path(r"C:\Users\mohan\Desktop") / "CyberSatDetect_Poster_Academic.pptx"
    prs.save(out)
    print("Saved:", out)
    return out


if __name__ == "__main__":
    build()
