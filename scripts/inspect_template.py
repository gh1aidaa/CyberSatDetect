"""Inspect the poster pptx template to understand its layout."""
import sys
from pptx import Presentation
from pptx.util import Emu

TEMPLATE = r"C:\Users\mohan\Desktop\_\u20681 - \u0642\u0627\u0644\u0628 \u0645\u0644\u062e\u0635 \u0627\u0644\u0628\u062d\u062b (1)\u2069 (3).pptx"

import glob, os
matches = glob.glob(r"C:\Users\mohan\Desktop\*.pptx")
for m in matches:
    print("FOUND:", repr(m))

# Use the third pptx
target = None
for m in matches:
    if "_" in os.path.basename(m) and "(3)" in m:
        target = m
        break

if target is None:
    print("Template not found, exiting")
    sys.exit(1)

print("\nUsing template:", repr(target))
prs = Presentation(target)
print("Slide width:", prs.slide_width, "EMU =", Emu(prs.slide_width).cm, "cm")
print("Slide height:", prs.slide_height, "EMU =", Emu(prs.slide_height).cm, "cm")
print("Number of slides:", len(prs.slides))

for si, slide in enumerate(prs.slides):
    print(f"\n=== Slide {si} ===")
    print("Layout:", slide.slide_layout.name)
    for shi, shape in enumerate(slide.shapes):
        print(f"  Shape {shi}: name={shape.name!r} type={shape.shape_type} "
              f"pos=({Emu(shape.left or 0).cm:.2f},{Emu(shape.top or 0).cm:.2f})cm "
              f"size=({Emu(shape.width or 0).cm:.2f}x{Emu(shape.height or 0).cm:.2f})cm")
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    text = run.text.replace("\n", "\\n")
                    print(f"    p{pi}r{ri}: {text!r}  font={run.font.name} size={run.font.size}")
