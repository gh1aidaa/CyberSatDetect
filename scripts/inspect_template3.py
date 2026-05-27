"""Dump the XML for key shapes so we can see Arabic characters correctly."""
import glob, os, sys
from pptx import Presentation
from lxml import etree

matches = glob.glob(r"C:\Users\mohan\Desktop\*.pptx")
target = None
for m in matches:
    if "(3)" in m and "_" in os.path.basename(m):
        target = m
        break

prs = Presentation(target)
slide = prs.slides[0]

# Print XML for shapes of interest
target_indices = [1, 3, 5, 11, 13, 20, 23, 24, 25, 26, 27, 28]
out_dir = r"C:\Users\mohan\Desktop\gh\CyberSatDetectprojct2\CyberSatDetectprojct\scripts\template_xml"
os.makedirs(out_dir, exist_ok=True)

for i, shape in enumerate(slide.shapes):
    if i not in target_indices:
        continue
    xml_str = etree.tostring(shape._element, pretty_print=True, encoding="unicode")
    fp = os.path.join(out_dir, f"shape_{i:02d}_{shape.name.replace(' ', '_')}.xml")
    try:
        with open(fp, "w", encoding="utf-8") as f:
            f.write(xml_str)
        print(f"Wrote {fp}")
    except Exception as e:
        print(f"Failed to write {i}: {e}")

# Print layout master
layout_xml = etree.tostring(slide.slide_layout._element, pretty_print=True, encoding="unicode")
with open(os.path.join(out_dir, "_layout.xml"), "w", encoding="utf-8") as f:
    f.write(layout_xml)
print("Wrote layout xml")
