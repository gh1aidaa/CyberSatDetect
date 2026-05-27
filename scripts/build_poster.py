# -*- coding: utf-8 -*-
"""
Build the official CyberSatDetect graduation poster (A0, 84.1 x 118.9 cm) by
filling the provided template `_⁨1 - قالب ملخص البحث (1)⁩ (3).pptx`.

Content sources (all from inside this repository):
  - docs/pipeline_report.md
  - docs/CSD_continual_learning_guide_AR.md
  - thesis_official_evaluation_figures/evaluation_summary_4attacks.json
  - thesis_official_evaluation_figures/overall_threshold_metrics_4attacks.csv
  - thesis_official_evaluation_figures/per_attack_full_metrics_4attacks.csv

Nothing in this script is invented — every number/term is taken from the
project files. Section labels are aligned with the official poster guidelines
(Introduction, Methodology, Tools, Results, Conclusion, Recommendations).
"""
from __future__ import annotations

import copy
import glob
import io
import os
import sys

# Force UTF-8 on stdout so Arabic/Unicode prints work in the Windows console.
if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")
from typing import List

from pptx import Presentation
from pptx.util import Emu, Pt
from lxml import etree

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _find_template() -> str:
    matches = glob.glob(r"C:\Users\mohan\Desktop\*.pptx")
    for m in matches:
        base = os.path.basename(m)
        if base.startswith("_") and "(3)" in base:
            return m
    raise FileNotFoundError("Template _⁨1 - قالب ملخص البحث (1)⁩ (3).pptx not found")


# -----------------------------------------------------------------------------
# Low level helpers
# -----------------------------------------------------------------------------

def _clear_paragraphs(tx_body: etree._Element) -> None:
    """Remove all <a:p> children from a text body, keep <a:bodyPr> and <a:lstStyle>."""
    for p in tx_body.findall(f"{{{NS['a']}}}p"):
        tx_body.remove(p)


def _make_run(rPr_template: etree._Element | None, text: str, *,
              size_pt: int | None = None, bold: bool | None = None,
              color_hex: str | None = None) -> etree._Element:
    """Create an <a:r> element using a copy of rPr_template for formatting."""
    r = etree.SubElement(etree.Element(f"{{{NS['a']}}}__tmp"),
                         f"{{{NS['a']}}}r")
    if rPr_template is not None:
        rPr = copy.deepcopy(rPr_template)
    else:
        rPr = etree.SubElement(r, f"{{{NS['a']}}}rPr")
        rPr.set("lang", "ar-SA")
        rPr.set("dirty", "0")
    if size_pt is not None:
        rPr.set("sz", str(int(size_pt * 100)))
    if bold is not None:
        rPr.set("b", "1" if bold else "0")
    if color_hex is not None:
        # Replace any fill on rPr with a solid fill of the requested colour.
        for fill in list(rPr):
            tag = etree.QName(fill).localname
            if tag in {"solidFill", "noFill", "gradFill", "blipFill", "pattFill"}:
                rPr.remove(fill)
        solid = etree.SubElement(rPr, f"{{{NS['a']}}}solidFill")
        srgb = etree.SubElement(solid, f"{{{NS['a']}}}srgbClr")
        srgb.set("val", color_hex)
    r.append(rPr)
    t = etree.SubElement(r, f"{{{NS['a']}}}t")
    t.text = text
    return r


def _new_paragraph(rtl: bool = True, align: str = "just",
                   bullet: bool = False, line_spacing_pts: int | None = None) -> etree._Element:
    p = etree.Element(f"{{{NS['a']}}}p")
    pPr = etree.SubElement(p, f"{{{NS['a']}}}pPr")
    pPr.set("marL", "0")
    pPr.set("marR", "0")
    pPr.set("lvl", "0")
    pPr.set("indent", "0")
    pPr.set("algn", align)
    pPr.set("defTabSz", "457200")
    if rtl:
        pPr.set("rtl", "1")
    pPr.set("eaLnBrk", "1")
    pPr.set("fontAlgn", "auto")
    pPr.set("latinLnBrk", "0")
    pPr.set("hangingPunct", "1")
    if line_spacing_pts is not None:
        lnSpc = etree.SubElement(pPr, f"{{{NS['a']}}}lnSpc")
        spcPts = etree.SubElement(lnSpc, f"{{{NS['a']}}}spcPts")
        spcPts.set("val", str(line_spacing_pts * 100))
    if bullet:
        buFont = etree.SubElement(pPr, f"{{{NS['a']}}}buFont")
        buFont.set("typeface", "Arial")
        buChar = etree.SubElement(pPr, f"{{{NS['a']}}}buChar")
        buChar.set("char", "•")
    else:
        etree.SubElement(pPr, f"{{{NS['a']}}}buNone")
    return p


def _grab_rPr_template(shape) -> etree._Element | None:
    """Pull the first <a:rPr> in a shape's text body, useful as formatting source."""
    body = shape.text_frame._txBody  # noqa: SLF001
    rPr = body.find(f".//{{{NS['a']}}}rPr")
    return rPr


def write_lines(shape, lines: List[str], *, size_pt: int = 22,
                bold_first: bool = False, color_hex: str | None = None,
                align: str = "just", line_spacing_pts: int = 30,
                bullet: bool = False) -> None:
    """Replace the shape's text body with a list of paragraph strings.

    Each entry becomes one paragraph. The first paragraph can be bold (for
    section intros). Formatting is inherited from the existing template
    runs where possible.
    """
    body = shape.text_frame._txBody  # noqa: SLF001
    rPr_tpl = _grab_rPr_template(shape)
    _clear_paragraphs(body)
    for i, line in enumerate(lines):
        p = _new_paragraph(rtl=True, align=align, bullet=bullet,
                           line_spacing_pts=line_spacing_pts)
        run = _make_run(rPr_tpl, line, size_pt=size_pt,
                        bold=(bold_first and i == 0) or None,
                        color_hex=color_hex)
        p.append(run)
        body.append(p)


def write_single(shape, text: str, *, size_pt: int | None = None,
                 bold: bool | None = None, color_hex: str | None = None,
                 align: str = "ctr") -> None:
    """Replace the shape's text body with a single paragraph/run."""
    body = shape.text_frame._txBody  # noqa: SLF001
    rPr_tpl = _grab_rPr_template(shape)
    _clear_paragraphs(body)
    p = _new_paragraph(rtl=True, align=align)
    run = _make_run(rPr_tpl, text, size_pt=size_pt, bold=bold, color_hex=color_hex)
    p.append(run)
    body.append(p)


# -----------------------------------------------------------------------------
# Content (all sourced from this repository)
# -----------------------------------------------------------------------------

PROJECT_TITLE_AR = "CyberSatDetect"
PROJECT_SUBTITLE_AR = "نظام كشف الشذوذ المدعوم بالذكاء الاصطناعي لأمن الأقمار الصناعية"

# Right-column banner 1 — "Authors"
RIGHT_HEADER_1_AR = "إعداد الطالبات  /  Project Team"
# Right-column banner 2 — "College / Affiliation"
RIGHT_HEADER_2_AR = "الجامعة والكلية  /  Affiliation"

# Middle column section headers (Arabic; English labels already in template)
SEC_INTRO_AR     = "مقدمة"
SEC_AIMS_AR      = "أهداف البحث"
SEC_METHODS_AR   = "منهجية وطرق البحث"
SEC_RESULTS_AR   = "النتائج والمناقشة"
SEC_RECOMM_AR    = "التوصيات"

INTRO_PARS = [
    "تُعدّ الأقمار الصناعية بنية تحتية حيوية لخدمات الاتصالات والملاحة ومراقبة الأرض، وأي خلل في بيانات قياسها عن بُعد (Telemetry) قد يهدد المهام ويعرّض الأصول الفضائية للخطر.",
    "تستهدف الهجمات السيبرانية هذه البيانات بأنماط متعددة (تجميد، طفرات، انحراف، ضوضاء) قد لا تكشفها القواعد الثابتة، مما يستوجب نظام كشف يتعلّم السلوك الطبيعي ويرصد الانحرافات.",
    "يقدّم مشروع CyberSatDetect نظام كشف شذوذ غير مُشرَف يدرّب نموذجاً هجيناً (Autoencoder + Predictor) على بيانات طبيعية فقط، ويزوّده بوحدة تعلم مستمر آمنة للتكيّف مع توزيعات تشغيل جديدة دون نسيان كارثي.",
]

AIMS_PARS = [
    "• بناء نموذج هجين (Autoencoder + Predictor) يتعلّم من البيانات الطبيعية فقط لكشف الشذوذ في تيليمتري الأقمار الصناعية.",
    "• تحقيق فصل عالي الجودة بين السلوك الطبيعي وأربعة أنماط هجوم: Drift, Freeze, Noise, Spike.",
    "• تقليل الإنذارات الكاذبة عبر عتبات إحصائية مشتقّة من البيانات الطبيعية فقط (p95, p97, p99, p99.5, p99.7, 3σ).",
    "• تمكين تعلّم مستمر آمن: انتقاء نوافذ طبيعية موثوقة من بيئة التشغيل، ودمجها بضبط دقيق (fine-tuning) مع replay للحماية من النسيان.",
    "• توفير حوكمة تشغيلية: سجل نماذج (model_registry) بمراحل PENDING / APPROVED مع تراجع آمن (rollback).",
]

METHODS_PARS = [
    "١) ابتلاع البيانات: قبول صيغتي CSV و NPY مع تنظيف رقمي صارم (float32، استبدال القيم غير المنتهية، استيفاء خطي، تعبئة الفجوات بالأصفار).",
    "٢) تجزئة بنوافذ انزلاقية: طول W = 100 وخطوة S = 50 (تداخل 50%)، وتنسيق المدخل إلى X ∈ ℝ^(B×100×1).",
    "٣) النموذج الهجين: رأس إعادة بناء + رأس تنبؤ بالخطوة التالية، مع خسارة مركّبة:  L_total = W_recon·L_recon + W_pred·L_pred + W_grad·L_grad + W_sep·L_sep.",
    "٤) درجة الشذوذ وقت التشغيل (بدون حد الفصل):  score = e_recon + e_pred + e_grad.",
    "٥) عتبات إحصائية تُحسب من البيانات الطبيعية فقط: p99 / p99.5 / p99.7 / 3σ.",
    "٦) تقييم صارم على بيانات attacked_v2 — تُسمَّى النافذة شاذة إذا تجاوزت نسبة خطواتها المهاجَمة 10%.",
    "الأدوات: Python • TensorFlow / Keras • NumPy • FastAPI • React • Docker.",
]

RESULTS_PARS = [
    "بروتوكول QC-filtered • مجموعة الاختبار: 72 ملفاً • 226,623 نافذة طبيعية • 4 أنماط هجوم (drift, freeze, noise, spike) — 54,712 نافذة شاذة من 182,323.",
    "عتبة best-F1:  F1 = 0.9493 • Accuracy = 98.57% • Balanced Accuracy = 99.14% • Recall = 0.9991 • Precision = 0.9043 • FAR ≈ 1.63%.",
    "منحنيات الأداء:  ROC-AUC = 0.9958 • PR-AUC = 0.9561.",
    "عتبة إحصائية p99 (إنذارات منخفضة):  Recall = 0.7539 • FAR = 0.47% • F1 = 0.8451.",
    "الخلاصة: فصل قوي بين الطبيعي والمهاجَم، وتوازن مرن بين الاستدعاء و FAR عبر اختيار العتبة المناسبة للتشغيل.",
]

RECOMM_PARS = [
    "الخلاصة:  أثبت CyberSatDetect أن نموذجاً هجيناً (Autoencoder + Predictor) مدرَّباً على بيانات طبيعية فقط يكشف أربعة أنماط هجوم رئيسية على تيليمتري الأقمار الصناعية بأداء عالٍ (F1 = 0.949 و ROC-AUC = 0.996) مع FAR ≈ 1.63%.",
    "",
    "التوصيات العملية:",
    "• نشر النظام داخل المحطات الأرضية كطبقة مراقبة في الوقت الفعلي، وتكامله مع منظومات SIEM / SOC.",
    "• توسيع تغطية الهجمات إلى أنماط مركّبة (Drift+Spike، Pattern Shift، Scale، Drop) لرفع التعميم.",
    "• ضبط عتبة التشغيل حسب ميزانية FAR: p99 لخفض الإنذارات، و best-F1 لرفع الاستدعاء عند الحاجة.",
    "• تفعيل دورة التعلم المستمر بسياسة موافقة بشرية صارمة:  build_dataset ← fine_tune ← model_registry (PENDING) ← موافقة الأدمن ← APPROVED.",
    "• تخزين سجلات تفسيرية لكل نافذة شاذة (مكوّنات e_recon, e_pred, e_grad) لدعم تحليلات ما بعد الحادث.",
    "• إجراء تقييم دوري على بيانات إنتاج جديدة لرصد انجراف التوزيع وضبط العتبات وفقاً للسلوك المُلاحَظ.",
    "",
    "أعمال مستقبلية:  دعم تيليمتري متعدد القنوات • تقييم على بيئة تشغيلية حيّة • سياسة عتبة ديناميكية تتكيّف مع انجراف التوزيع • إدماج هجمات مركّبة • توسعة التعلم المستمر إلى تعلم متعدد المهام.",
]

# Right column box 1 — Project Team (Shape 36)
TEAM_PARS = [
    "كلية الحاسبات — قسم الأمن السيبراني — جامعة أم القرى",
    "",
    "إعداد الطالبات:",
    "• غيداء بنت علي القرني",
    "• آلاء بنت فهد الحازمي",
    "• بيان بنت علي الزهراني",
    "• غيداء بنت عبدالعزيز المسعودي",
    "",
    "تحت إشراف:  د. عهد بنت عبدالرحمن الجرف",
]

# Right column box 2 — Affiliation / Project info (Shape 38)
AFFIL_PARS = [
    "جامعة أم القرى",
    "كلية الحاسبات • قسم الأمن السيبراني",
    "",
    "مشروع التخرج — ٢٠٢٦",
    "اسم المشروع:  CyberSatDetect",
    "نظام كشف الشذوذ المدعوم بالذكاء الاصطناعي لأمن الأقمار الصناعية.",
    "",
    "نطاق المشروع:  بناء، تدريب، تقييم، ونشر نموذج هجين لكشف الشذوذ في تيليمتري الأقمار الصناعية، مع وحدة تعلم مستمر آمنة وحوكمة بنماذج معتمَدة.",
]

TEAM_LINE = "تواصل المشروع:  قسم الأمن السيبراني — كلية الحاسبات — جامعة أم القرى"
SUPERVISOR_LINE = "المشروع البرمجي:  CyberSatDetect • Python • TensorFlow • FastAPI • React"


# -----------------------------------------------------------------------------
# Main
# -----------------------------------------------------------------------------

def main() -> int:
    template_path = _find_template()
    out_path = os.path.join(
        r"C:\Users\mohan\Desktop",
        "CyberSatDetect_Poster_A0_Final.pptx",
    )
    print(f"Template: {template_path}")
    print(f"Output:   {out_path}")

    prs = Presentation(template_path)
    slide = prs.slides[0]
    shapes = list(slide.shapes)

    # ------------------------------------------------------------------
    # Title (Shape 28) — top banner spanning the full width
    # ------------------------------------------------------------------
    title_shape = shapes[28]
    body = title_shape.text_frame._txBody  # noqa: SLF001
    rPr_tpl = _grab_rPr_template(title_shape)
    _clear_paragraphs(body)

    p1 = _new_paragraph(rtl=True, align="ctr", line_spacing_pts=72)
    p1.append(_make_run(rPr_tpl, PROJECT_TITLE_AR, size_pt=64, bold=True,
                        color_hex="0B5566"))
    body.append(p1)

    p2 = _new_paragraph(rtl=True, align="ctr", line_spacing_pts=46)
    p2.append(_make_run(rPr_tpl, PROJECT_SUBTITLE_AR, size_pt=36, bold=True,
                        color_hex="1F3D55"))
    body.append(p2)

    # ------------------------------------------------------------------
    # Middle-column Arabic section headers
    # ------------------------------------------------------------------
    section_header_size = 40
    section_header_color = "0B5566"

    # Shape 1 (top middle - paired with English "Introduction")
    write_single(shapes[1], SEC_INTRO_AR, size_pt=section_header_size,
                 bold=True, color_hex=section_header_color, align="r")
    # Shape 3 (paired with "Aims")
    write_single(shapes[3], SEC_AIMS_AR, size_pt=section_header_size,
                 bold=True, color_hex=section_header_color, align="r")
    # Shape 5 (paired with "Methods") — was broken into single-character runs
    write_single(shapes[5], SEC_METHODS_AR, size_pt=section_header_size,
                 bold=True, color_hex=section_header_color, align="r")
    # Shape 20 (paired with "Results & Discussions")
    write_single(shapes[20], SEC_RESULTS_AR, size_pt=section_header_size,
                 bold=True, color_hex=section_header_color, align="r")
    # Shape 23 (paired with "Recommendation")
    write_single(shapes[23], SEC_RECOMM_AR, size_pt=section_header_size,
                 bold=True, color_hex=section_header_color, align="r")

    # ------------------------------------------------------------------
    # Right-column group banners (Shape 11, 13) — text inside the group
    # ------------------------------------------------------------------
    def _replace_group_banner_text(group_shape, new_text: str) -> None:
        # The text element inside the group has name "نص 5".
        for sub in group_shape.shapes:
            if sub.has_text_frame:
                rPr_tpl = _grab_rPr_template(sub)
                body = sub.text_frame._txBody  # noqa: SLF001
                _clear_paragraphs(body)
                p = _new_paragraph(rtl=True, align="ctr", line_spacing_pts=46)
                p.append(_make_run(rPr_tpl, new_text, size_pt=34, bold=True,
                                   color_hex="FFFFFF"))
                body.append(p)

    _replace_group_banner_text(shapes[11], RIGHT_HEADER_1_AR)
    _replace_group_banner_text(shapes[13], RIGHT_HEADER_2_AR)

    # ------------------------------------------------------------------
    # Body content
    # ------------------------------------------------------------------
    body_size = 22
    body_color = "1F3D55"

    write_lines(shapes[31], INTRO_PARS, size_pt=body_size,
                color_hex=body_color, line_spacing_pts=30)
    write_lines(shapes[30], AIMS_PARS, size_pt=body_size,
                color_hex=body_color, line_spacing_pts=30)
    write_lines(shapes[14], METHODS_PARS, size_pt=body_size,
                color_hex=body_color, line_spacing_pts=30)
    write_lines(shapes[33], RESULTS_PARS, size_pt=body_size,
                color_hex=body_color, line_spacing_pts=30)
    write_lines(shapes[32], RECOMM_PARS, size_pt=body_size,
                color_hex=body_color, line_spacing_pts=32)

    write_lines(shapes[36], TEAM_PARS, size_pt=body_size,
                color_hex=body_color, line_spacing_pts=32)
    write_lines(shapes[38], AFFIL_PARS, size_pt=body_size,
                color_hex=body_color, line_spacing_pts=32)

    # ------------------------------------------------------------------
    # Bottom-right team / supervisor info (Shape 8 = bottom, Shape 9 = top)
    # ------------------------------------------------------------------
    write_single(shapes[9], TEAM_LINE, size_pt=18, bold=True,
                 color_hex=body_color, align="r")
    write_single(shapes[8], SUPERVISOR_LINE, size_pt=20, bold=True,
                 color_hex="0B5566", align="r")

    # Optional: clear the small Shape 12 textbox (was empty placeholder)
    # — leave untouched to preserve background graphics.

    prs.save(out_path)
    print("Saved poster:", out_path)
    return 0


if __name__ == "__main__":
    sys.exit(main())
