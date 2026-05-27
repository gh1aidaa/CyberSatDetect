"""Inspect the CyberSatDetect poster PPTX: list slides, shapes, positions, sizes, text."""
from pptx import Presentation
from pptx.util import Emu

SRC = r"C:\Users\mohan\Desktop\3.pptx"
OUT = "poster3_inspect.txt"


def emu_to_cm(v):
    return round(v / 360000, 2)


def main():
    import io
    out = io.StringIO()
    prs = Presentation(SRC)
    out.write(f"Slide size: {emu_to_cm(prs.slide_width)} x {emu_to_cm(prs.slide_height)} cm\n")
    out.write(f"Num slides: {len(prs.slides)}\n")
    for si, slide in enumerate(prs.slides):
        out.write(f"\n=== Slide {si+1} ===\n")
        for i, shape in enumerate(slide.shapes):
            try:
                left = emu_to_cm(shape.left) if shape.left is not None else None
                top = emu_to_cm(shape.top) if shape.top is not None else None
                w = emu_to_cm(shape.width) if shape.width is not None else None
                h = emu_to_cm(shape.height) if shape.height is not None else None
            except Exception:
                left = top = w = h = None
            text = ""
            try:
                if shape.has_text_frame:
                    text = " | ".join(p.text for p in shape.text_frame.paragraphs)
            except Exception:
                pass
            text_short = (text[:200] + "...") if len(text) > 200 else text
            out.write(f"  [{i}] {shape.shape_type} name='{shape.name}' L={left} T={top} W={w} H={h}\n")
            if text_short.strip():
                out.write(f"      text: {text_short!r}\n")
    with open(OUT, "w", encoding="utf-8") as f:
        f.write(out.getvalue())
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
